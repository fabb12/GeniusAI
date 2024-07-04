import anthropic
from PyQt6.QtWidgets import QFileDialog, QMessageBox, QDialog, QVBoxLayout, QLabel, QPushButton, QHBoxLayout
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
#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

antrophic_key = "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"
model_3_5_sonnet = "claude-3-5-sonnet-20240620"
model_3_haiku = "claude-3-haiku-20240307"
class PptxGeneration:
    @staticmethod
    def impostaFont(shape, size_pt, text):
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.size = Pt(size_pt)

    @staticmethod
    def createPresentationFromFile(parent, file_path, num_slide, company_name, language):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                testo = file.read()
            testo_per_slide, input_tokens, output_tokens = PptxGeneration.generaTestoPerSlide(testo, num_slide,
                                                                                              company_name, language)
            print(f"Token di input utilizzati: {input_tokens}")
            print(f"Token di output utilizzati: {output_tokens}")

            save_path, _ = QFileDialog.getSaveFileName(parent, "Salva Presentazione", "",
                                                       "PowerPoint Presentation (*.pptx)")
            if save_path:
                PptxGeneration.createPresentationFromText(parent, testo_per_slide, save_path)
            else:
                QMessageBox.warning(parent, "Attenzione", "Salvataggio annullato. Nessun file selezionato.")
        except Exception as e:
            QMessageBox.critical(parent, "Errore", f"Si è verificato un errore durante la lettura del file: {e}")

    @staticmethod
    def generaTestoPerSlide(testo, num_slide, company_name, language):
        client = anthropic.Anthropic(api_key=antrophic_key)

        # Costruisci la parte del messaggio relativa alla compagnia se company_name è fornito
        company_info = (
            f" The presentation is targeted at the company {company_name}. "
            f"Gather all relevant information about {company_name}, including its scope, main products, "
            f"and the market it operates in. Use this information to create a personalized AI-generated presentation "
            f"for {company_name} based on the topic provided."
            if company_name else ""
        )

        message = client.messages.create(
            model=model_3_5_sonnet,
            max_tokens=1000,
            temperature=0.7,
            system=(
                    "You are a professional slide deck designer. Your task is to transform the following "
                    "text into a format suitable for PowerPoint slides. Each slide should include a title, "
                    "subtitle, and content with bullet points. "
                    f"Please follow this structure for {num_slide} slides:\n\n"
                    "Titolo: [Title of the slide]\n"
                    "Sottotitolo: [Subtitle of the slide]\n"
                    "Contenuto:\n- Bullet point 1\n- Bullet point 2\n\n"
                    f"Ensure the final presentation is in {language}." +
                    company_info
            ),
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"{testo}\n\nAssistant:"
                        }
                    ]
                }
            ]
        )
        # Estrai il testo risultante e i token utilizzati
        testo_resultante = message.content[0].text
        input_tokens = message.usage.input_tokens
        output_tokens = message.usage.output_tokens

        return testo_resultante, input_tokens, output_tokens

    @staticmethod
    def creaPresentazione(parent, transcriptionTextArea, num_slide, company_name, language):
        testo_attuale = transcriptionTextArea.toPlainText()
        if testo_attuale.strip() == "":
            file_path, _ = QFileDialog.getOpenFileName(parent, "Seleziona File di Testo", "", "Text Files (*.txt)")
            if file_path:
                PptxGeneration.createPresentationFromFile(parent, file_path)
            else:
                QMessageBox.warning(parent, "Attenzione", "Nessun testo inserito e nessun file selezionato.")
        else:
            save_path, _ = QFileDialog.getSaveFileName(parent, "Salva Presentazione", "",
                                                       "PowerPoint Presentation (*.pptx)")
            if save_path:
                testo_per_slide, input_tokens, output_tokens = PptxGeneration.generaTestoPerSlide(testo_attuale,
                                                                                                  num_slide,
                                                                                                  company_name,
                                                                                                  language)
                print(f"Token di input utilizzati: {input_tokens}")
                print(f"Token di output utilizzati: {output_tokens}")
                transcriptionTextArea.setPlainText(testo_per_slide)
                PptxGeneration.createPresentationFromText(parent, testo_per_slide, save_path)
            else:
                QMessageBox.warning(parent, "Attenzione", "Salvataggio annullato. Nessun file selezionato.")

    @staticmethod
    def createPresentationFromText(parent, testo, output_file):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]

        def imposta_testo_e_font(paragraph, text, size_pt, bold=False, color=None):
            run = paragraph.add_run()
            run.text = text
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(size_pt)
            font.bold = bold
            if color:
                font.color.rgb = RGBColor(*color)

        def aggiungi_paragrafo_formattato(shape, text, size_pt, bold=False, color=None, bullet=False):
            paragraph = shape.text_frame.add_paragraph()
            if bullet:
                paragraph.level = 0
            imposta_testo_e_font(paragraph, text, size_pt, bold, color)
            return paragraph

        def aggiungi_footer(slide):
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
            footer_text = footer.text_frame.add_paragraph()
            imposta_testo_e_font(footer_text, "Made by GeniusAI", 10, color=(255, 255, 255))
            footer_text.alignment = 2  # Centro allineato

        clean_text = re.sub(r'\*\*(Titolo|Sottotitolo|Contenuto):', r'\1:', testo)
        clean_text = re.sub(r'-\s*', '', clean_text)

        pattern = r"Titolo:\s*(.*?)\s+Sottotitolo:\s*(.*?)\s+Contenuto:\s*(.*?)\s*(?=Titolo|$)"
        slides_data = re.findall(pattern, clean_text, re.DOTALL)

        for index, (titolo_text, sottotitolo_text, contenuto_text) in enumerate(slides_data):
            if index == 0:
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                aggiungi_paragrafo_formattato(title, titolo_text.strip(), 44, bold=True, color=(0, 0, 0))
                aggiungi_paragrafo_formattato(subtitle, sottotitolo_text.strip(), 32, color=(89, 89, 89))
            else:
                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                aggiungi_paragrafo_formattato(title, titolo_text.strip(), 44, bold=True, color=(0, 0, 0))

                subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
                aggiungi_paragrafo_formattato(subtitle_shape, sottotitolo_text.strip(), 28, color=(89, 89, 89))

                for line in contenuto_text.strip().split('\n'):
                    line = line.strip()
                    if line.startswith('•'):
                        aggiungi_paragrafo_formattato(content, line[1:].strip(), 24, bullet=True)
                    elif ':' in line:
                        parts = line.split(':', 1)
                        p = aggiungi_paragrafo_formattato(content, parts[0].strip() + ':', 24, bold=True)
                        if len(parts) > 1:
                            imposta_testo_e_font(p, parts[1].strip(), 24)
                    else:
                        aggiungi_paragrafo_formattato(content, line, 24)

            aggiungi_footer(slide)  # Aggiunge il footer a ogni slide

        if prs.slides:
            prs.save(output_file)
            QMessageBox.information(None, "Successo",
                                    "Presentazione PowerPoint generata con successo e salvata in: " + output_file)
            PptxGeneration.visualizzaPresentazione(parent, output_file)

        else:               QMessageBox.warning(None, "Attenzione",
                                "Non sono state generate slides a causa di dati di input non validi o mancanti.")
    def visualizzaPresentazione(parent, file_path):
        prs = Presentation(file_path)
        dialog = QDialog(parent)
        dialog.setWindowTitle("Visualizza Presentazione")
        layout = QVBoxLayout()

        for i, slide in enumerate(prs.slides):
            slide_label = QLabel(f"Slide {i+1}: {slide.shapes.title.text if slide.shapes.title else 'Senza titolo'}")
            layout.addWidget(slide_label)

        buttonLayout = QHBoxLayout()
        fineSlideButton = QPushButton("Imposta Fine Slide")
        fineSlideButton.clicked.connect(lambda: PptxGeneration.impostaFineSlide(prs, dialog))
        buttonLayout.addWidget(fineSlideButton)

        layout.addLayout(buttonLayout)
        dialog.setLayout(layout)
        dialog.exec()

    @staticmethod
    def impostaFineSlide(prs, dialog):
        dialog.accept()
        QMessageBox.information(dialog, "Informazione", "Fine della slide impostata con successo.")
