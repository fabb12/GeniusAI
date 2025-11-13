# File: src/services/PptxGeneration.py

import anthropic
import google.generativeai as genai
import requests
import json
import logging
from PyQt6.QtWidgets import QFileDialog, QMessageBox, QDialog, QVBoxLayout, QLabel, QPushButton, QHBoxLayout, QApplication
from PyQt6.QtCore import Qt
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
import os
import subprocess  # Needed for os.startfile alternative
import sys
import tempfile
from dotenv import load_dotenv

# Importa la configurazione delle azioni e le chiavi/endpoint necessari
from src.config import (
    OLLAMA_ENDPOINT, get_api_key, get_model_for_action,
    PROMPT_PPTX_GENERATION # Assicurati che questo percorso sia corretto e il file esista
)
from src.services.utils import _call_ollama_api

load_dotenv() # Carica .env se necessario

class PptxGeneration:
    """
    Classe statica per gestire la generazione di presentazioni PowerPoint
    utilizzando diversi modelli AI (Claude, Gemini, Ollama) basati sulle impostazioni.
    """

    @staticmethod
    def _find_layout_by_placeholder_types(prs, *placeholder_types):
        """
        Trova un layout di slide che contenga tutti i tipi di placeholder specificati.
        Cerca il layout più "pulito" (con meno placeholder extra).
        """
        best_layout = None
        min_extra_placeholders = float('inf')

        for layout in prs.slide_layouts:
            layout_placeholder_types = {p.placeholder_format.type for p in layout.placeholders}

            # Verifica se il layout contiene tutti i placeholder richiesti
            if all(pt in layout_placeholder_types for pt in placeholder_types):
                num_extra_placeholders = len(layout_placeholder_types) - len(placeholder_types)

                # Preferisce layout con meno placeholder extra
                if num_extra_placeholders < min_extra_placeholders:
                    min_extra_placeholders = num_extra_placeholders
                    best_layout = layout
                    # Se troviamo una corrispondenza esatta, è l'ideale
                    if num_extra_placeholders == 0:
                        return best_layout

        if best_layout:
            logging.info(f"Trovato layout '{best_layout.name}' per i placeholder: {placeholder_types}")
        else:
            logging.warning(f"Nessun layout trovato per i placeholder: {placeholder_types}")

        return best_layout


    @staticmethod
    def createPresentationFromFile(parent, file_path, num_slide, company_name, language):
        """Crea una presentazione leggendo il testo sorgente da un file."""
        try:
            logging.info(f"Tentativo di creare presentazione da file: {file_path}")
            with open(file_path, 'r', encoding='utf-8') as file:
                testo = file.read()

            if not testo.strip():
                if hasattr(parent, 'show_status_message'):
                    parent.show_status_message("Il file selezionato è vuoto.", error=True)
                return

            result_tuple = PptxGeneration.generaTestoPerSlide(testo, num_slide, company_name, language)

            if not isinstance(result_tuple, tuple) or len(result_tuple) != 3:
                 QMessageBox.critical(parent, "Errore API", f"Errore durante la generazione del testo per le slide.\nDettagli:\n{result_tuple}")
                 return

            testo_per_slide, input_tokens, output_tokens = result_tuple
            logging.info(f"Generazione testo PPTX da file - Token Input: {input_tokens}, Token Output: {output_tokens}")

            save_path, _ = QFileDialog.getSaveFileName(parent, "Salva Presentazione", "", "PowerPoint Presentation (*.pptx)")
            if save_path:
                PptxGeneration.createPresentationFromText(parent, testo_per_slide, save_path, num_slides=num_slide)
            else:
                if hasattr(parent, 'show_status_message'):
                    parent.show_status_message("Salvataggio annullato.", error=True)
        except FileNotFoundError:
             logging.error(f"File non trovato: {file_path}")
             QMessageBox.critical(parent, "Errore", f"File non trovato: {file_path}")
        except Exception as e:
            logging.exception("Errore in createPresentationFromFile")
            QMessageBox.critical(parent, "Errore", f"Si è verificato un errore imprevisto: {e}")

    @staticmethod
    def generaTestoPerSlide(testo, num_slide, company_name, language):
        """
        Genera il testo strutturato per le slide usando il modello AI selezionato.
        Restituisce una tupla (testo_generato, input_tokens, output_tokens) o una stringa di errore.
        """
        # Recupera il modello selezionato per l'azione 'pptx_generation'
        selected_model = get_model_for_action('pptx_generation')
        logging.info(f"Generazione PPTX: Utilizzo del modello '{selected_model}'")

        # 2. Leggi e formatta il prompt
        try:
            with open(PROMPT_PPTX_GENERATION, 'r', encoding='utf-8') as f:
                prompt_template = f.read()
        except Exception as e:
            logging.exception("Errore durante la lettura del file prompt PPTX")
            return f"Errore lettura prompt: {e}"

        company_info = (
            f" La presentazione è destinata all'azienda {company_name}. "
            f"Considera il suo ambito, i prodotti principali e il mercato. "
            f"Crea una presentazione personalizzata per {company_name} basata sull'argomento fornito."
            if company_name else ""
        )
        try:
            system_prompt_content = prompt_template.format(
                num_slide=num_slide,
                language=language,
                company_info=company_info
            )
        except KeyError as e:
             logging.error(f"Errore nella formattazione del prompt PPTX. Chiave mancante: {e}")
             return f"Errore nel template del prompt: chiave '{e}' mancante."

        user_prompt = f"Testo sorgente:\n{testo}\n\n---\nGenera la struttura della presentazione come richiesto." # Input principale per l'LLM

        # 3. Selezione e chiamata API
        model_name_lower = selected_model.lower()

        try:
            if "ollama:" in model_name_lower:
                logging.info(f"Chiamata API Ollama: {selected_model}")
                ollama_model_name = selected_model.split(":", 1)[1]

                result_text = _call_ollama_api(
                    OLLAMA_ENDPOINT,
                    ollama_model_name,
                    system_prompt_content,
                    user_prompt
                )
                logging.info("Risposta ricevuta da Ollama.")
                return result_text, 0, 0

            elif "gemini" in model_name_lower:
                # --- Logica per Google Gemini ---
                logging.info(f"Chiamata API Gemini: {selected_model}")
                google_api_key = get_api_key('google')
                if not google_api_key: raise ValueError("API Key Google non configurata.")
                genai.configure(api_key=google_api_key)
                model = genai.GenerativeModel(selected_model, system_instruction=system_prompt_content) # Usa system instruction
                response = model.generate_content(user_prompt) # Passa solo lo user prompt
                # Aggiungi gestione errori risposta Gemini se necessario (es. blocco per safety)
                try:
                    result_text = response.text
                except ValueError as ve:
                     # Potrebbe essere bloccato per safety
                     logging.warning(f"Possibile blocco risposta Gemini: {ve}. Dettagli: {response.prompt_feedback}")
                     raise Exception(f"Risposta da Gemini bloccata o non valida. Causa: {response.prompt_feedback}")

                logging.info("Risposta ricevuta da Gemini.")
                return result_text, 0, 0 # Token non disponibili facilmente

            elif "claude" in model_name_lower:
                # --- Logica per Anthropic Claude ---
                logging.info(f"Chiamata API Claude: {selected_model}")
                anthropic_api_key = get_api_key('anthropic')
                if not anthropic_api_key: raise ValueError("API Key Anthropic non configurata.")
                client = anthropic.Anthropic(api_key=anthropic_api_key)
                message = client.messages.create(
                    model=selected_model,
                    max_tokens=4096,
                    temperature=0.7,
                    system=system_prompt_content,
                    messages=[{"role": "user", "content": [{"type": "text", "text": user_prompt}]}]
                )
                # Aggiungi controllo sul motivo di stop
                if message.stop_reason == 'max_tokens':
                     logging.warning("Risposta Claude troncata per max_tokens.")

                testo_resultante = message.content[0].text
                input_tokens = message.usage.input_tokens
                output_tokens = message.usage.output_tokens
                logging.info("Risposta ricevuta da Claude.")
                return testo_resultante, input_tokens, output_tokens

            # elif "gpt" in model_name_lower:
                # --- Logica per OpenAI (Esempio) ---
                # logging.info(f"Chiamata API OpenAI: {selected_model}")
                # if not OPENAI_API_KEY: raise ValueError("API Key OpenAI non configurata.")
                # from openai import OpenAI
                # client = OpenAI(api_key=OPENAI_API_KEY)
                # response = client.chat.completions.create(
                #     model=selected_model,
                #     messages=[
                #         {"role": "system", "content": system_prompt_content},
                #         {"role": "user", "content": user_prompt}
                #     ],
                #     temperature=0.7,
                #     max_tokens=4000
                # )
                # testo_resultante = response.choices[0].message.content
                # input_tokens = response.usage.prompt_tokens
                # output_tokens = response.usage.completion_tokens
                # logging.info("Risposta ricevuta da OpenAI.")
                # return testo_resultante, input_tokens, output_tokens

            else:
                logging.error(f"Modello '{selected_model}' non gestito per PPTX.")
                return f"Errore: Modello '{selected_model}' non supportato."

        except requests.exceptions.ConnectionError:
             logging.error(f"Impossibile connettersi a Ollama: {OLLAMA_ENDPOINT}")
             return f"Errore di connessione a Ollama ({OLLAMA_ENDPOINT}). Verifica che sia in esecuzione."
        except requests.exceptions.Timeout:
             logging.error(f"Timeout durante connessione a Ollama ({selected_model})")
             return f"Timeout Ollama ({selected_model}). Il modello potrebbe essere lento o non rispondere."
        except Exception as e:
             logging.exception(f"Errore API durante la generazione del testo per PPTX con {selected_model}")
             return f"Errore API ({type(e).__name__}): {str(e)}"


    @staticmethod
    def creaPresentazione(parent, transcriptionTextArea, num_slide, company_name, language, template_path=None):
        """Metodo principale chiamato dall'UI per creare la presentazione."""
        testo_attuale = transcriptionTextArea.toPlainText()
        if not testo_attuale.strip():
            if hasattr(parent, 'show_status_message'):
                parent.show_status_message("Inserisci del testo prima di generare.", error=True)
            return

        save_path, _ = QFileDialog.getSaveFileName(parent, "Salva Presentazione Come", "", "PowerPoint Presentation (*.pptx)")
        if not save_path:
            if hasattr(parent, 'show_status_message'):
                parent.show_status_message("Salvataggio annullato.", error=True)
            return

        # Avvia il thread per la generazione
        parent.start_pptx_generation_thread(
            testo_attuale, save_path, template_path, num_slide, company_name, language
        )


    @staticmethod
    def _find_placeholder(slide, placeholder_type):
        """Trova il primo placeholder di un tipo specifico."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == placeholder_type:
                return shape
        return None

    @staticmethod
    def _find_content_placeholder(slide):
        """Trova il placeholder principale per il contenuto, cercando BODY o OBJECT."""
        logging.info(f"Cerco placeholder BODY in '{slide.slide_layout.name}'")
        placeholders = [p for p in slide.placeholders if p.placeholder_format.type == PP_PLACEHOLDER.BODY]
        if placeholders:
            logging.info(f"Trovati {len(placeholders)} placeholder BODY.")
            return max(placeholders, key=lambda p: p.width * p.height) if len(placeholders) > 1 else placeholders[0]

        logging.info(f"Nessun placeholder BODY trovato. Cerco placeholder OBJECT.")
        placeholders = [p for p in slide.placeholders if p.placeholder_format.type == PP_PLACEHOLDER.OBJECT]
        if placeholders:
            logging.info(f"Trovati {len(placeholders)} placeholder OBJECT.")
            return max(placeholders, key=lambda p: p.width * p.height) if len(placeholders) > 1 else placeholders[0]

        logging.warning("Nessun placeholder BODY o OBJECT trovato.")
        return None

    @staticmethod
    def _add_content_to_slide(slide, content_text):
        """Aggiunge il contenuto principale (bullet points) a una slide, gestendo i livelli di indentazione."""
        body_shape = PptxGeneration._find_content_placeholder(slide)
        if not body_shape:
            logging.warning(f"Layout '{slide.slide_layout.name}' non ha un placeholder per il contenuto (BODY/OBJECT). Salto il contenuto.")
            return

        tf = body_shape.text_frame
        tf.clear()

        # Rimuove l'eventuale paragrafo vuoto creato da clear()
        if len(tf.paragraphs) > 0 and not tf.paragraphs[0].text.strip():
            p = tf.paragraphs[0]
            p._element.getparent().remove(p._element)

        for line in content_text.strip().split('\n'):
            if not line.strip():
                continue

            p = tf.add_paragraph()
            p.text = re.sub(r'^\s*[-*•]\s*', '', line).strip()
            p.level = (len(line) - len(line.lstrip(' '))) // 2

    @staticmethod
    def _truncate_slides(slides_data, num_slides):
        """Tronca o avvisa se il numero di slide generate non corrisponde a quello richiesto."""
        if num_slides is None:
            return slides_data

        if len(slides_data) > num_slides:
            logging.warning(f"L'AI ha generato {len(slides_data)} slide, ma ne sono state richieste {num_slides}. Tronco le slide extra.")
            return slides_data[:num_slides]
        elif len(slides_data) < num_slides:
            logging.warning(f"L'AI ha generato {len(slides_data)} slide, ma ne sono state richieste {num_slides}. Non verranno aggiunte slide vuote.")
        return slides_data

    @staticmethod
    def createPresentationFromText(parent, testo, output_file, template_path=None, num_slides=None, is_preview=False):
        """Crea il file .pptx dal testo strutturato generato dall'AI, usando un approccio robusto per i layout."""
        logging.info(f"Tentativo di creare file PPTX: {output_file} con template: {template_path}")
        try:
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
                # Rimuovi tutte le slide esistenti dal template
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            else:
                prs = Presentation()

            title_slide_layout = PptxGeneration._find_layout_by_placeholder_types(prs, PP_PLACEHOLDER.TITLE)
            content_slide_layout = PptxGeneration._find_layout_by_placeholder_types(prs, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY)

            if not content_slide_layout:
                if parent:
                    QMessageBox.critical(parent, "Errore Template", "Impossibile trovare un layout 'Titolo e Contenuto' adeguato nel template.")
                return
            if not title_slide_layout:
                logging.warning("Nessun layout solo titolo trovato, userò 'Titolo e Contenuto' come fallback.")
                title_slide_layout = content_slide_layout

            clean_text = re.sub(r'\*\*(Titolo|Sottotitolo|Contenuto)\*\*:', r'\1:', testo, flags=re.IGNORECASE)
            slide_blocks = re.split(r'\n(?=Titolo:)', clean_text.strip(), flags=re.IGNORECASE)

            slides_data = []
            for block in slide_blocks:
                if not block.strip():
                    continue

                titolo_match = re.search(r'Titolo:\s*(.*)', block, re.IGNORECASE)
                sottotitolo_match = re.search(r'Sottotitolo:\s*(.*)', block, re.IGNORECASE)
                contenuto_match = re.search(r'Contenuto:\s*((.|\n)*)', block, re.IGNORECASE)

                titolo = titolo_match.group(1).strip() if titolo_match else ''
                sottotitolo = sottotitolo_match.group(1).strip() if sottotitolo_match else ''
                contenuto = contenuto_match.group(1).strip() if contenuto_match else ''

                if titolo:
                    slides_data.append({
                        'titolo': titolo,
                        'sottotitolo': sottotitolo,
                        'contenuto': contenuto
                    })

            if not slides_data:
                if hasattr(parent, 'show_status_message'):
                    parent.show_status_message("Impossibile estrarre dati strutturati dal testo dell'AI.", error=True)
                return

            if num_slides is not None:
                slides_data = PptxGeneration._truncate_slides(slides_data, num_slides)

            if not slides_data:
                return

            first_slide_data = slides_data.pop(0)
            layout_per_prima_slide = content_slide_layout if first_slide_data.get('contenuto', '').strip() else title_slide_layout
            slide = prs.slides.add_slide(layout_per_prima_slide)

            if slide.shapes.title:
                slide.shapes.title.text = first_slide_data.get('titolo', '').strip()
            subtitle_shape = PptxGeneration._find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
            if subtitle_shape and first_slide_data.get('sottotitolo'):
                subtitle_shape.text = first_slide_data.get('sottotitolo', '').strip()
            if first_slide_data.get('contenuto'):
                PptxGeneration._add_content_to_slide(slide, first_slide_data.get('contenuto', '').strip())

            for slide_data in slides_data:
                slide = prs.slides.add_slide(content_slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get('titolo', '').strip()
                subtitle_shape = PptxGeneration._find_placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
                if subtitle_shape and slide_data.get('sottotitolo'):
                    subtitle_shape.text = slide_data.get('sottotitolo', '').strip()
                if slide_data.get('contenuto'):
                    PptxGeneration._add_content_to_slide(slide, slide_data.get('contenuto', '').strip())

            if prs.slides:
                prs.save(output_file)
                logging.info(f"Presentazione salvata con successo: {output_file}")
                # La notifica e l'apertura del file sono gestite dal thread
            else:
                logging.warning("Nessuna slide è stata generata.")

        except Exception as e:
            logging.exception("Errore durante la creazione del file PPTX")
            raise e # Rilancia l'eccezione per essere catturata dal thread

    @staticmethod
    def generate_preview(parent, testo, template_path=None, num_slides=None):
        """Genera un'anteprima della presentazione come immagini."""
        if sys.platform != "win32":
            QMessageBox.warning(parent, "Funzionalità non supportata", "La generazione dell'anteprima è supportata solo su Windows.")
            return None

        try:
            import win32com.client
        except ImportError:
            QMessageBox.critical(parent, "Dipendenza Mancante", "La libreria 'pywin32' è necessaria per l'anteprima. Installala con 'pip install pywin32'.")
            return None

        temp_pptx_file = None
        temp_image_dir = parent.get_temp_dir()
        image_paths = []

        try:
            # 1. Crea un file .pptx temporaneo
            temp_pptx_file = parent.get_temp_filepath(suffix=".pptx")

            # Passa num_slides per coerenza con la generazione finale e is_preview per disabilitare il popup
            PptxGeneration.createPresentationFromText(
                parent, testo, temp_pptx_file, template_path, num_slides=num_slides, is_preview=True
            )

            # 2. Usa COM per esportare le slide come immagini
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(temp_pptx_file, WithWindow=False)

            for i, slide in enumerate(presentation.Slides):
                image_path = os.path.join(temp_image_dir, f"slide_{i + 1}.png")
                slide.Export(image_path, "PNG")
                image_paths.append(image_path)

            presentation.Close()
            powerpoint.Quit()

            return image_paths

        except Exception as e:
            logging.exception("Errore durante la generazione dell'anteprima")
            QMessageBox.critical(parent, "Errore Anteprima", f"Si è verificato un errore: {e}")
            return None
        finally:
            # Pulizia file temporaneo
            if temp_pptx_file and os.path.exists(temp_pptx_file):
                os.remove(temp_pptx_file)
