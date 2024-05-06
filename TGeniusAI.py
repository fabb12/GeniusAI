import sys
from pydub import AudioSegment
from PyQt6.QtGui import QAction
from PyQt6.QtWidgets import (QFileDialog,  QMessageBox,QSizePolicy)
from PyQt6.QtCore import QUrl
from moviepy.editor import  concatenate_videoclips, concatenate_audioclips
from moviepy.editor import VideoFileClip, vfx, AudioFileClip, AudioClip, ImageClip, CompositeVideoClip
from moviepy.editor import  ColorClip
from pptx import Presentation
import re
import tempfile
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QPushButton, QLabel, QCheckBox,
                             QLineEdit, QSlider, QHBoxLayout, QGroupBox, QTextEdit, QRadioButton, QComboBox, QButtonGroup)
from PyQt6.QtGui import QIcon, QFont
from PyQt6.QtMultimediaWidgets import QVideoWidget
from PyQt6.QtMultimedia import QMediaPlayer, QAudioOutput
from PyQt6.QtCore import Qt, QSize
from pyqtgraph.dockarea.Dock import Dock
from pyqtgraph.dockarea.DockArea import DockArea
import datetime
import os
import time
from PyQt6.QtWidgets import QProgressDialog
from DownloadVideo import DownloadThread
from AudioTranscript import TranscriptionThread
from AudioGenerationREST import AudioGenerationThread
from VideoCutting import VideoCuttingThread
import cv2
from ScreenRecorder import ScreenRecorder
import pygetwindow as gw
from screeninfo import get_monitors
from SettingsManager import DockSettingsManager
import sounddevice as sd
from PyQt6.QtCore import QTimer, QTime
from pptx.util import Pt, Inches
import speech_recognition as sr
from num2words import num2words
from langdetect import detect, LangDetectException
import pycountry
import uuid
from CropVideo import CropVideoWidget
from moviepy.audio.AudioClip import CompositeAudioClip
from PyQt6.QtCore import pyqtSignal

class VideoAudioManager(QMainWindow):
    def __init__(self):
        super().__init__()

        #self.setGeometry(100, 500, 800, 800)
        self.player = QMediaPlayer()
        self.audioOutput = QAudioOutput()  # Crea un'istanza di QAudioOutput
        self.playerOutput = QMediaPlayer()
        self.audioOutputOutput = QAudioOutput()

        self.player.setAudioOutput(self.audioOutput)  # Imposta l'audio output del player
        self.audioOutput.setVolume(1.0)  # Imposta il volume al massimo (1.0 = 100%)
        self.recentFiles = []
        self.initUI()
        self.currentPosition = 0
        self.videoPathLineEdit = ''
        self.videoPathLineOutputEdit = ''
        self.is_recording = False
        self.video_writer = None
        self.setupDockSettingsManager()
        self.current_video_path = None
        self.current_audio_path = None

    def initUI(self):

        self.setWindowTitle('ThemaGeniusAI - Alpha')
        self.setWindowIcon(QIcon('res/eye.png'))

        # Creazione e configurazione dell'area del dock
        area = DockArea()
        self.setCentralWidget(area)

        # Creazione dei docks esistenti...
        self.videoPlayerDock = Dock("Video Player Source")
        self.videoPlayerDock.setStyleSheet(self.styleSheet())

        self.transcriptionDock = Dock("Trascrizione e Sintesi Audio")
        self.transcriptionDock.setStyleSheet(self.styleSheet())

        self.editingDock = Dock("Opzioni di Editing")
        self.editingDock.setStyleSheet(self.styleSheet())
        self.downloadDock = self.createDownloadDock()
        self.downloadDock.setStyleSheet(self.styleSheet())
        self.recordingDock = self.createRecordingDock()
        self.recordingDock.setStyleSheet(self.styleSheet())
        self.audioDock = self.createAudioDock()  # Creazione del dock per la gestione audio
        self.audioDock.setStyleSheet(self.styleSheet())
        self.videoPlayerOutputDock = Dock("Video Player Output")
        self.videoPlayerOutputDock.setStyleSheet(self.styleSheet())

        # Creazione del dock merge videos
        self.videoMergeDock = self.createVideoMergeDock()
        self.videoMergeDock.setStyleSheet(self.styleSheet())

        # Aggiunta dei docks all'area
        area.addDock(self.videoPlayerOutputDock, 'right')  # Posizionamento a destra
        area.addDock(self.audioDock, 'bottom')  # Aggiungi il dock audio alla posizione desiderata
        area.addDock(self.videoPlayerDock, 'left')
        area.addDock(self.transcriptionDock, 'bottom')
        area.addDock(self.editingDock, 'right')
        area.addDock(self.downloadDock, 'left')
        area.addDock(self.recordingDock, 'right')
        area.addDock(self.videoMergeDock, 'bottom')

        if hasattr(self, 'applyDarkMode'):
            self.applyDarkMode()

        self.applyStyleToAllDocks()  # Applica lo stile dark a tutti i dock
        # Setup del dock del video player
        self.videoCropWidget = CropVideoWidget()
        self.videoCropWidget.setAcceptDrops(True)
        self.videoCropWidget.setSizePolicy(QSizePolicy.Policy.Expanding,
                                        QSizePolicy.Policy.Expanding)
        self.player.setVideoOutput(self.videoCropWidget)


        self.videoSlider = QSlider(Qt.Orientation.Horizontal)

        # Label per mostrare il nome del file video
        self.fileNameLabel = QLabel("Nessun video caricato")
        self.fileNameLabel.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.fileNameLabel.setStyleSheet("QLabel { font-weight: bold; }")

        # Creazione dei pulsanti di controllo playback
        self.playButton = QPushButton('')
        self.playButton.setIcon(QIcon("res/play.png"))
        self.pauseButton = QPushButton('')
        self.pauseButton.setIcon(QIcon("res/pausa.png"))
        self.stopButton = QPushButton('')
        self.stopButton.setIcon(QIcon("res/stop.png"))
        self.cutButton = QPushButton('')  # Se necessario
        self.cutButton.setIcon(QIcon("res/taglia.png"))
        self.cropButton = QPushButton('Ritaglia')  # Se necessario

        # Collegamento dei pulsanti ai loro slot funzionali
        self.playButton.clicked.connect(self.playVideo)
        self.pauseButton.clicked.connect(self.pauseVideo)
        self.stopButton.clicked.connect(self.stopVideo)
        self.cutButton.clicked.connect(self.cutVideo)  # Assumendo che la funzione cutVideo sia definita
        self.cropButton.clicked.connect(self.applyCrop)  # Assumendo che la funzione cutVideo sia definita

        # Creazione e configurazione del display del timecode
        self.currentTimeLabel = QLabel('00:00')
        self.totalTimeLabel = QLabel('/ 00:00')
        timecodeLayout = QHBoxLayout()
        timecodeLayout.addWidget(self.currentTimeLabel)
        timecodeLayout.addWidget(self.totalTimeLabel)


        # Video Player output
        # Setup del widget video per l'output
        videoOutputWidget = QVideoWidget()
        videoOutputWidget.setAcceptDrops(True)
        videoOutputWidget.setSizePolicy(QSizePolicy.Policy.Expanding,
                                        QSizePolicy.Policy.Expanding)

        self.playerOutput.setAudioOutput(self.audioOutputOutput)
        self.playerOutput.setVideoOutput(videoOutputWidget)

        # Creazione dei pulsanti di controllo playback per il video output
        playButtonOutput = QPushButton('')
        playButtonOutput.setIcon(QIcon("res/play.png"))
        pauseButtonOutput = QPushButton('')
        pauseButtonOutput.setIcon(QIcon("res/pausa.png"))
        stopButtonOutput = QPushButton('')
        stopButtonOutput.setIcon(QIcon("res/stop.png"))

        changeButtonOutput = QPushButton('')
        changeButtonOutput.setIcon(QIcon("res/change.png"))
        changeButtonOutput.setToolTip('Sposta video in Video Player Source')
        changeButtonOutput.clicked.connect(lambda: self.loadVideo(self.videoPathLineOutputEdit,
                                                                  os.path.basename(self.videoPathLineOutputEdit)))

        # Collegamento dei pulsanti ai loro slot funzionali
        playButtonOutput.clicked.connect(lambda: self.playerOutput.play())
        pauseButtonOutput.clicked.connect(lambda: self.playerOutput.pause())
        stopButtonOutput.clicked.connect(lambda: self.playerOutput.stop())

        # Layout per i controlli di playback
        playbackControlLayoutOutput = QHBoxLayout()
        playbackControlLayoutOutput.addWidget(playButtonOutput)
        playbackControlLayoutOutput.addWidget(pauseButtonOutput)
        playbackControlLayoutOutput.addWidget(stopButtonOutput)
        playbackControlLayoutOutput.addWidget(changeButtonOutput)

        # Slider per il controllo della posizione del video output
        videoSliderOutput = QSlider(Qt.Orientation.Horizontal)
        videoSliderOutput.setRange(0, 1000)  # Inizializza con un range di esempio
        videoSliderOutput.sliderMoved.connect(lambda position: self.playerOutput.setPosition(position))


        # Creazione delle QLabel per il timecode
        self.currentTimeLabelOutput = QLabel('00:00')
        self.totalTimeLabelOutput = QLabel('/ 00:00')
        timecodeLayoutOutput = QHBoxLayout()
        timecodeLayoutOutput.addWidget(self.currentTimeLabelOutput)
        timecodeLayoutOutput.addWidget(self.totalTimeLabelOutput)

        # Inserisci il layout del timecode nel layout principale del video output


        # Label per mostrare il nome del file video output
        self.fileNameLabelOutput = QLabel("Nessun video caricato")
        self.fileNameLabelOutput.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.fileNameLabelOutput.setStyleSheet("QLabel { font-weight: bold; }")

        # Layout principale per il dock del video player output
        videoOutputLayout = QVBoxLayout()
        videoOutputLayout.addWidget(self.fileNameLabelOutput)

        videoOutputLayout.addWidget(videoOutputWidget)
        videoOutputLayout.addLayout(timecodeLayoutOutput)
        videoOutputLayout.addWidget(videoSliderOutput)
        videoOutputLayout.addLayout(playbackControlLayoutOutput)


        self.playerOutput.durationChanged.connect(self.updateDurationOutput)
        self.playerOutput.positionChanged.connect(self.updateTimeCodeOutput)

        # Widget per contenere il layout del video player output
        videoPlayerOutputWidget = QWidget()
        videoPlayerOutputWidget.setLayout(videoOutputLayout)
        self.videoPlayerOutputDock.addWidget(videoPlayerOutputWidget)

        # Collegamento degli eventi del player multimediale ai metodi corrispondenti
        self.playerOutput.durationChanged.connect(lambda duration: videoSliderOutput.setRange(0, duration))
        self.playerOutput.positionChanged.connect(lambda position: videoSliderOutput.setValue(position))

        # trascrizione video
        self.transcribeButton = QPushButton('Trascrivi Video')
        self.transcribeButton.clicked.connect(self.transcribeVideo)

        # Layout per i controlli di playback
        playbackControlLayout = QHBoxLayout()
        playbackControlLayout.addWidget(self.playButton)
        playbackControlLayout.addWidget(self.pauseButton)
        playbackControlLayout.addWidget(self.stopButton)
        playbackControlLayout.addWidget(self.cutButton)
        #playbackControlLayout.addWidget(self.cropButton)

        # Layout principale per il dock del video player
        videoPlayerLayout = QVBoxLayout()
        videoPlayerLayout.addWidget(self.fileNameLabel)
        videoPlayerLayout.addWidget(self.videoCropWidget)  # Aggiunta del widget video
        videoPlayerLayout.addLayout(timecodeLayout)  # Aggiunta del display del timecode
        videoPlayerLayout.addWidget(self.videoSlider)  # Aggiunta della slider

        videoPlayerLayout.addLayout(playbackControlLayout)  # Aggiunta dei controlli di playback
        videoPlayerLayout.addWidget(self.transcribeButton)  # Aggiunta della slider

        # Widget per contenere il layout del video player
        videoPlayerWidget = QWidget()
        videoPlayerWidget.setLayout(videoPlayerLayout)
        self.videoPlayerDock.addWidget(videoPlayerWidget)

        # Setup del dock di trascrizione e sintesi audio
        transGroupBox = QGroupBox("Gestione Trascrizione")

        # Creazione di un layout interno per il GroupBox
        innerLayout = QVBoxLayout()

        # Layout orizzontale per i pulsanti "Incolla" e "Salva"
        buttonsLayout = QHBoxLayout()

        # Inizializzazione della QLabel per la lingua della trascrizione
        self.transcriptionLanguageLabel = QLabel("Lingua rilevata: Nessuna")
        self.transcriptionLanguageLabel.setAlignment(Qt.AlignmentFlag.AlignLeft)

        # Layout orizzontale per la label e la combo box della selezione della lingua
        languageSelectionLayout = QHBoxLayout()
        languageLabel = QLabel("Seleziona lingua video:")  # Creazione della QLabel per il testo
        languageLabel.setAlignment(Qt.AlignmentFlag.AlignLeft)  # Allinea il testo a destra

        # Inizializzazione della QComboBox per la lingua
        self.languageComboBox = QComboBox()

        self.languageComboBox.addItem("Italiano", "it")
        self.languageComboBox.addItem("Inglese", "en")
        self.languageComboBox.addItem("Francese", "fr")
        self.languageComboBox.addItem("Spagnolo", "es")
        self.languageComboBox.addItem("Tedesco", "de")
        self.languageComboBox.currentIndexChanged.connect(
            self.onLanguageChange)  # Opzionale: gestire il cambio di lingua

        # Aggiunta della label e della combo box al layout orizzontale
        languageSelectionLayout.addWidget(languageLabel)
        languageSelectionLayout.addWidget(self.languageComboBox)
        languageSelectionLayout.addStretch(1)
        # Aggiunta del layout di selezione della lingua al layout interno del GroupBox
        innerLayout.addLayout(languageSelectionLayout)  # Usa addLayout qui

        # TextArea per la trascrizione
        self.transcriptionTextArea = CustomTextEdit(self)
        self.transcriptionTextArea.setStyleSheet("""
               QTextEdit {
                   color: white;
                   font-size: 12pt;
                   font-family: 'Arial';
                   background-color: #333;
               }
           """)
        self.transcriptionTextArea.setPlaceholderText("Incolla qui la tua trascrizione...")
        self.transcriptionTextArea.textChanged.connect(self.handleTextChange)
        self.resetButton = QPushButton()
        self.resetButton.setIcon(QIcon("res/reset.png"))  # Assicurati che il percorso dell'icona sia corretto
        self.resetButton.setFixedSize(24, 24)  # Imposta la dimensione del pulsante
        self.resetButton.clicked.connect(lambda: self.transcriptionTextArea.clear())
        self.detected_language_code = 'it-IT'  # Imposta una lingua di default
        self.video_download_language = None
        # Pulsante per incollare nel QTextEdit
        self.pasteButton = QPushButton()
        self.pasteButton.setIcon(QIcon("res/paste.png"))  # Assicurati che il percorso dell'icona sia corretto
        self.pasteButton.setFixedSize(24, 24)  # Imposta la dimensione del pulsante
        self.pasteButton.clicked.connect(lambda: self.transcriptionTextArea.paste())

        # Pulsante per salvare il testo
        self.saveButton = QPushButton()
        self.saveButton.setIcon(QIcon("res/save.png"))  # Assicurati che il percorso dell'icona sia corretto
        self.saveButton.setFixedSize(24, 24)  # Imposta la dimensione del pulsante
        self.saveButton.clicked.connect(self.saveText)

        # Aggiungi i pulsanti "Incolla" e "Salva" al layout orizzontale
        buttonsLayout.addWidget(self.resetButton)
        buttonsLayout.addWidget(self.pasteButton)
        buttonsLayout.addWidget(self.saveButton)
        buttonsLayout.setAlignment(Qt.AlignmentFlag.AlignLeft)

        # Pulsanti per le diverse funzionalità
        self.generateAudioButton = QPushButton('Genera Audio con AI')
        self.generateAudioButton.clicked.connect(self.generateAudioWithElevenLabs)

        self.generatePresentationButton = QPushButton('Genera Presentazione')
        self.generatePresentationButton.clicked.connect(self.generaPresentationConTestoAttuale)
        # Aggiunta dei layout e widget al layout interno
        innerLayout.addLayout(buttonsLayout)  # Aggiungi il layout dei pulsanti in orizzontale
        innerLayout.addWidget(self.transcriptionTextArea)
        innerLayout.addWidget(self.transcriptionLanguageLabel)
        bottonLayout = QHBoxLayout()
        bottonLayout.addWidget(self.generateAudioButton)
        bottonLayout.addWidget(self.generatePresentationButton)
        innerLayout.addLayout(bottonLayout)

        # Impostazione del layout interno al GroupBox
        transGroupBox.setLayout(innerLayout)

        # Creazione del widget e assegnazione del layout con GroupBox
        widgetTranscription = QWidget()
        widgetTranscription.setLayout(QVBoxLayout())
        widgetTranscription.layout().addWidget(transGroupBox)

        # Aggiunta del widget al dock
        self.transcriptionDock.addWidget(widgetTranscription)
        # Layout principale per il dock delle opzioni di editing

        voiceSettingsWidget = self.setupVoiceSettingsUI()
        self.editingDock.addWidget(voiceSettingsWidget)

        # Setup della barra dei menu e della dark mode, se necessario
        self.setupMenuBar()


        # Collegamento degli eventi del player multimediale ai metodi corrispondenti
        self.player.durationChanged.connect(self.durationChanged)  # Assicurati che questo slot sia definito
        self.player.positionChanged.connect(self.positionChanged)  # Assicurati che questo slot sia definito

        self.videoSlider.sliderMoved.connect(self.setPosition)  # Assicurati che questo slot sia definito


    def keyPressEvent(self, event):
        if event.key() == Qt.Key.Key_Right:
            # Avanza di 2 secondi
            current_position = self.player.position()
            new_position = current_position + 2000  # 2000 ms = 2 secondi
            self.player.setPosition(new_position)
        elif event.key() == Qt.Key.Key_Left:
            # Torna indietro di 2 secondi
            current_position = self.player.position()
            new_position = max(0, current_position - 2000)  # Evita di andare sotto lo 0
            self.player.setPosition(new_position)
        elif event.key() == Qt.Key.Key_Space:
            # Pausa o riproduzione a seconda dello stato corrente
            if self.player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
                self.player.pause()
            else:
                self.player.play()
        else:
            super().keyPressEvent(event)  # gestione degli altri eventi di tastiera

    def updateTimeCodeOutput(self, position):
        # Aggiorna il timecode corrente del video output
        hours, remainder = divmod(position // 1000, 3600)
        minutes, seconds = divmod(remainder, 60)
        self.currentTimeLabelOutput.setText(f'{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}')

    def updateDurationOutput(self, duration):
        # Aggiorna la durata totale del video output
        hours, remainder = divmod(duration // 1000, 3600)
        minutes, seconds = divmod(remainder, 60)
        self.totalTimeLabelOutput.setText(f' / {int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}')


    def applyCrop(self):
        if not self.videoPathLineEdit or not os.path.exists(self.videoPathLineEdit):
            QMessageBox.warning(self, "Errore", "Carica un video prima di applicare il ritaglio.")
            return

        cropRect = self.videoCropWidget.getCropRect()
        if cropRect.isEmpty():
            QMessageBox.warning(self, "Errore", "Seleziona un'area da ritagliare.")
            return

        try:
            video = VideoFileClip(self.videoPathLineEdit)
            cropped_video = video.crop(x1=cropRect.x(), y1=cropRect.y(), x2=cropRect.x() + cropRect.width(),
                                       y2=cropRect.y() + cropRect.height())
            output_path = tempfile.mktemp(suffix='.mp4')
            cropped_video.write_videofile(output_path, codec='libx264')
            QMessageBox.information(self, "Successo", f"Il video ritagliato è stato salvato in {output_path}")

            self.loadVideoOutput(output_path)
        except Exception as e:
            QMessageBox.critical(self, "Errore durante il ritaglio", str(e))

    def applyStyleToAllDocks(self):
        style = self.getDarkStyle()
        self.videoPlayerDock.setStyleSheet(style)
        self.transcriptionDock.setStyleSheet(style)
        self.editingDock.setStyleSheet(style)
        self.downloadDock.setStyleSheet(style)
        self.recordingDock.setStyleSheet(style)
        self.audioDock.setStyleSheet(style)
        self.videoPlayerOutputDock.setStyleSheet(style)
    def getDarkStyle(self):
        return """
        QWidget {
            background-color: #2b2b2b;
            color: #dcdcdc;
        }
        QPushButton {
            background-color: #555555;
            border: 1px solid #666666;
            border-radius: 2px;
            padding: 5px;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #666666;
        }
        QPushButton:pressed {
            background-color: #777777;
        }
        QLabel {
            color: #cccccc;
        }
        QLineEdit {
            background-color: #333333;
            border: 1px solid #555555;
            border-radius: 2px;
            padding: 5px;
            color: #ffffff;
        }
        """
    def setupVoiceSettingsUI(self):
        voiceSettingsGroup = QGroupBox("Impostazioni Voce")
        layout = QVBoxLayout()

        # QComboBox per la selezione della voce
        self.voiceSelectionComboBox = QComboBox()
        self.voiceSelectionComboBox.addItem("Marco", "GcAgjAjkhWsmUd4GlPiv")
        self.voiceSelectionComboBox.addItem("Alessio", "BTpQARcEj1XqVxdZjTI7")
        self.voiceSelectionComboBox.addItem("Matilda", "atq1BFi5ZHt88WgSOJRB")
        self.voiceSelectionComboBox.addItem("Mika", "B2j2knC2POvVW0XJE6Hi")
        layout.addWidget(self.voiceSelectionComboBox)

        # Radio buttons per la selezione del genere vocale

        # Slider per la stabilità
        stabilityLabel = QLabel("Stabilità:")
        self.stabilitySlider = QSlider(Qt.Orientation.Horizontal)
        self.stabilitySlider.setMinimum(0)
        self.stabilitySlider.setMaximum(100)
        self.stabilitySlider.setValue(50)  # Valore predefinito
        self.stabilitySlider.setToolTip(
            "Regola l'emozione e la coerenza. Minore per più emozione, maggiore per coerenza.")
        self.stabilityValueLabel = QLabel("50%")  # Visualizza il valore corrente
        self.stabilitySlider.valueChanged.connect(lambda value: self.stabilityValueLabel.setText(f"{value}%"))
        layout.addWidget(stabilityLabel)
        layout.addWidget(self.stabilitySlider)
        layout.addWidget(self.stabilityValueLabel)

        # Slider per la similarità
        similarityLabel = QLabel("Similarità:")
        self.similaritySlider = QSlider(Qt.Orientation.Horizontal)
        self.similaritySlider.setMinimum(0)
        self.similaritySlider.setMaximum(100)
        self.similaritySlider.setValue(80)  # Valore predefinito
        self.similaritySlider.setToolTip(
            "Determina quanto la voce AI si avvicina all'originale. Alti valori possono includere artefatti.")
        self.similarityValueLabel = QLabel("80%")  # Visualizza il valore corrente
        self.similaritySlider.valueChanged.connect(lambda value: self.similarityValueLabel.setText(f"{value}%"))
        layout.addWidget(similarityLabel)
        layout.addWidget(self.similaritySlider)
        layout.addWidget(self.similarityValueLabel)

        # Slider per l'esagerazione dello stile
        styleLabel = QLabel("Esagerazione Stile:")
        self.styleSlider = QSlider(Qt.Orientation.Horizontal)
        self.styleSlider.setMinimum(0)
        self.styleSlider.setMaximum(10)
        self.styleSlider.setValue(0)  # Valore predefinito
        self.styleSlider.setToolTip("Amplifica lo stile del parlante originale. Impostare a 0 per maggiore stabilità.")
        self.styleValueLabel = QLabel("0")  # Visualizza il valore corrente
        self.styleSlider.valueChanged.connect(lambda value: self.styleValueLabel.setText(f"{value}"))
        layout.addWidget(styleLabel)
        layout.addWidget(self.styleSlider)
        layout.addWidget(self.styleValueLabel)

        # Checkbox per l'uso di speaker boost
        self.speakerBoostCheckBox = QCheckBox("Usa Speaker Boost")
        self.speakerBoostCheckBox.setChecked(True)
        self.speakerBoostCheckBox.setToolTip(
            "Potenzia la somiglianza col parlante originale a costo di maggiori risorse.")
        layout.addWidget(self.speakerBoostCheckBox)

        voiceSettingsGroup.setLayout(layout)
        return voiceSettingsGroup


    def applyFreezeFramePause(self):
        video_path = self.videoPathLineEdit
        if not video_path or not os.path.exists(video_path):
            QMessageBox.warning(self, "Errore", "Carica un video prima di applicare una pausa.")
            return

        try:
            # Estrazione del timecode e della durata della pausa
            timecode = self.timecodePauseLineEdit.text()
            pause_duration = int(self.pauseDurationLineEdit.text())
            hours, minutes, seconds = map(int, timecode.split(':'))
            start_time = hours * 3600 + minutes * 60 + seconds

            # Caricamento del video
            video_clip = VideoFileClip(video_path)

            # Ottenere l'ultimo frame dal timecode specificato
            freeze_frame = video_clip.get_frame(start_time)

            # Creare un clip di immagine dal frame congelato e impostare la sua durata
            freeze_clip = ImageClip(freeze_frame).set_duration(pause_duration).set_fps(video_clip.fps)

            # Creazione di due parti di video originali
            original_video_part1 = video_clip.subclip(0, start_time)
            original_video_part2 = video_clip.subclip(start_time)

            # Combinazione delle clip video senza l'audio
            video_only = concatenate_videoclips(
                [original_video_part1.without_audio(), freeze_clip, original_video_part2.without_audio()],
                method="compose")

            # Ricreare il video con l'audio originale
            final_video = CompositeVideoClip([video_only.set_audio(video_clip.audio)])

            # Salvataggio del video finale
            output_path = tempfile.mktemp(suffix='.mp4')
            final_video.write_videofile(output_path, codec='libx264')
            QMessageBox.information(self, "Successo", f"Video con pausa frame congelato salvato in {output_path}")
            self.loadVideoOutput(output_path)
        except Exception as e:
            QMessageBox.critical(self, "Errore durante l'applicazione della pausa frame congelato", str(e))

    def createAudioDock(self):
        dock = Dock("Gestione Audio")
        layout = QVBoxLayout()

        # GroupBox per la sostituzione dell'audio principale
        audioReplacementGroup = self.createAudioReplacementGroup()
        layout.addWidget(audioReplacementGroup)

        # GroupBox per l'applicazione delle pause audio
        audioPauseGroup = self.createAudioPauseGroup()
        layout.addWidget(audioPauseGroup)

        # GroupBox per l'applicazione delle pause video
        videoPauseGroup = self.createVideoPauseGroup()
        layout.addWidget(videoPauseGroup)

        # GroupBox per la gestione dell'audio di sottofondo
        backgroundAudioGroup = self.createBackgroundAudioGroup()
        layout.addWidget(backgroundAudioGroup)

        widget = QWidget()
        widget.setLayout(layout)
        dock.addWidget(widget)

        return dock

    def createAudioReplacementGroup(self):
        audioReplacementGroup = QGroupBox("Sostituzione Audio Principale")
        layout = QVBoxLayout()

        self.audioPathLineEdit = QLineEdit()
        self.audioPathLineEdit.setReadOnly(True)
        browseAudioButton = QPushButton('Scegli Audio Principale')
        browseAudioButton.clicked.connect(self.browseAudio)
        applyAudioButton = QPushButton('Applica Audio Principale')
        applyAudioButton.clicked.connect(
            lambda: self.applyNewAudioToVideo(self.videoPathLineEdit, self.audioPathLineEdit.text()))

        layout.addWidget(self.audioPathLineEdit)
        layout.addWidget(browseAudioButton)
        layout.addWidget(applyAudioButton)
        audioReplacementGroup.setLayout(layout)

        return audioReplacementGroup

    def createAudioPauseGroup(self):
        audioPauseGroup = QGroupBox("Applica Pause Audio")
        layout = QVBoxLayout()

        self.pauseBeforeLineEdit = QLineEdit()
        self.pauseBeforeLineEdit.setPlaceholderText("Durata pausa iniziale (s)")
        self.pauseAfterLineEdit = QLineEdit()
        self.pauseAfterLineEdit.setPlaceholderText("Durata pausa finale (s)")
        applyPauseButton = QPushButton('Applica Pause Audio')
        applyPauseButton.clicked.connect(self.applyAudioWithPauses)

        layout.addWidget(QLabel("Pausa iniziale (s):"))
        layout.addWidget(self.pauseBeforeLineEdit)
        layout.addWidget(QLabel("Pausa finale (s):"))
        layout.addWidget(self.pauseAfterLineEdit)
        layout.addWidget(applyPauseButton)
        audioPauseGroup.setLayout(layout)

        return audioPauseGroup

    def createVideoPauseGroup(self):
        videoPauseGroup = QGroupBox("Applica Pausa Video")
        layout = QVBoxLayout()

        self.timecodePauseLineEdit = QLineEdit()
        self.timecodePauseLineEdit.setPlaceholderText("Inserisci Timecode (hh:mm:ss)")
        self.pauseDurationLineEdit = QLineEdit()
        self.pauseDurationLineEdit.setPlaceholderText("Durata Pausa (secondi)")
        applyVideoPauseButton = QPushButton('Applica Pausa Video')
        applyVideoPauseButton.clicked.connect(self.applyFreezeFramePause)

        layout.addWidget(QLabel("Timecode Inizio Pausa:"))
        layout.addWidget(self.timecodePauseLineEdit)
        layout.addWidget(QLabel("Durata Pausa (s):"))
        layout.addWidget(self.pauseDurationLineEdit)
        layout.addWidget(applyVideoPauseButton)
        videoPauseGroup.setLayout(layout)

        return videoPauseGroup

    def createBackgroundAudioGroup(self):
        backgroundAudioGroup = QGroupBox("Gestione Audio di Sottofondo")
        layout = QVBoxLayout()

        self.backgroundAudioPathLineEdit = QLineEdit()
        self.backgroundAudioPathLineEdit.setReadOnly(True)
        browseBackgroundAudioButton = QPushButton('Scegli Sottofondo')
        browseBackgroundAudioButton.clicked.connect(self.browseBackgroundAudio)
        self.volumeSlider = QSlider(Qt.Orientation.Horizontal)
        self.volumeSlider.setRange(0, 100)
        self.volumeSlider.setValue(50)
        self.volumeSlider.valueChanged.connect(self.adjustBackgroundVolume)
        applyBackgroundButton = QPushButton('Applica Sottofondo al Video')
        applyBackgroundButton.clicked.connect(self.applyBackgroundAudioToVideo)

        layout.addWidget(self.backgroundAudioPathLineEdit)
        layout.addWidget(browseBackgroundAudioButton)
        layout.addWidget(QLabel("Volume Sottofondo:"))
        layout.addWidget(self.volumeSlider)
        layout.addWidget(applyBackgroundButton)
        backgroundAudioGroup.setLayout(layout)

        return backgroundAudioGroup

    def createVideoMergeDock(self):
        """Crea e restituisce il dock per la gestione dell'unione di video."""
        dock = Dock("Unione Video")

        # GroupBox per organizzare visivamente le opzioni di unione video
        mergeGroup = QGroupBox("Opzioni di Unione Video")
        mergeLayout = QVBoxLayout()

        # Widget per selezionare il video da unire
        self.mergeVideoPathLineEdit = QLineEdit()
        self.mergeVideoPathLineEdit.setReadOnly(True)
        browseMergeVideoButton = QPushButton('Scegli Video da Unire')
        browseMergeVideoButton.clicked.connect(self.browseMergeVideo)

        # Widget per inserire il timecode
        self.timecodeLineEdit = QLineEdit()
        self.timecodeLineEdit.setPlaceholderText("Inserisci il timecode (formato hh:mm:ss)")

        # Pulsante per unire il video
        mergeButton = QPushButton('Unisci Video')
        mergeButton.clicked.connect(self.mergeVideo)

        # Aggiunta dei controlli al layout della GroupBox
        mergeLayout.addWidget(self.mergeVideoPathLineEdit)
        mergeLayout.addWidget(browseMergeVideoButton)
        mergeLayout.addWidget(self.timecodeLineEdit)
        mergeLayout.addWidget(mergeButton)

        # Imposta il layout del GroupBox
        mergeGroup.setLayout(mergeLayout)

        # Widget principale per il dock
        widget = QWidget()
        widget.setLayout(QVBoxLayout())
        widget.layout().addWidget(mergeGroup)
        dock.addWidget(widget)

        return dock
    def browseMergeVideo(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Seleziona Video da Unire", "",
                                                  "Video Files (*.mp4 *.mov *.avi)")
        if fileName:
            self.mergeVideoPathLineEdit.setText(fileName)

    def mergeVideo(self):
        base_video_path = self.videoPathLineEdit
        merge_video_path = self.mergeVideoPathLineEdit.text()
        timecode = self.timecodeLineEdit.text()

        if not base_video_path or not os.path.exists(base_video_path):
            QMessageBox.warning(self, "Errore", "Carica il video principale prima di unirne un altro.")
            return

        if not merge_video_path or not os.path.exists(merge_video_path):
            QMessageBox.warning(self, "Errore", "Seleziona un video da unire.")
            return

        try:
            base_clip = VideoFileClip(base_video_path)
            merge_clip = VideoFileClip(merge_video_path)
            # Converti il timecode in secondi
            tc_hours, tc_minutes, tc_seconds = map(int, timecode.split(':'))
            tc_seconds_total = tc_hours * 3600 + tc_minutes * 60 + tc_seconds

            # Unisci i video inserendo il secondo video al timecode specificato
            final_clip = concatenate_videoclips([
                base_clip.subclip(0, tc_seconds_total),
                merge_clip,
                base_clip.subclip(tc_seconds_total)
            ], method='compose')

            # Genera il percorso del file di output
            base_dir = os.path.dirname(base_video_path)
            base_name = os.path.splitext(os.path.basename(base_video_path))[0]
            output_path = os.path.join(base_dir, f"{base_name}_merged.mp4")

            final_clip.write_videofile(output_path, codec='libx264')
            QMessageBox.information(self, "Successo", f"Il video unito è stato salvato in {output_path}")

            # Aggiorna il percorso del video per il player
            self.loadVideoOutput(output_path)
        except Exception as e:
            QMessageBox.warning(self, "Errore", "Si è verificato un errore durante l'unione dei video: " + str(e))

    def browseBackgroundAudio(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Seleziona Audio di Sottofondo", "",
                                                  "Audio Files (*.mp3 *.wav)")
        if fileName:
            self.backgroundAudioPathLineEdit.setText(fileName)

    def adjustBackgroundVolume(self, value):
        # Qui dovrai implementare la logica per regolare il volume del sottofondo
        print(f"Volume del sottofondo regolato al {value}%")
    def setupDockSettingsManager(self):
        docks = {
            'videoPlayerDock': self.videoPlayerDock,
            'transcriptionDock': self.transcriptionDock,
            'editingDock': self.editingDock,
            'downloadDock': self.downloadDock,
            'recordingDock': self.recordingDock,
            'audioDock': self.audioDock,
            'videoPlayerDockOutput': self.videoPlayerOutputDock,
            'videoMergeDock': self.videoMergeDock
        }
        self.dockSettingsManager = DockSettingsManager(self, docks)

        self.dockSettingsManager.load_settings()


    def closeEvent(self, event):
        self.dockSettingsManager.save_settings()
        event.accept()

    def createRecordingDock(self):
        """Crea e restituisce il dock per la gestione della registrazione dello schermo."""
        dock = Dock("Registrazione")

        self.timer = QTimer(self)
        self.timer.timeout.connect(self.updateTimecodeRec)
        self.timecodeLabel = QLabel('00:00')
        self.timecodeLabel.setAlignment(Qt.AlignmentFlag.AlignLeft)

        self.recordingStatusLabel = QLabel("Stato: Pronto per la registrazione")
        self.recordingStatusLabel.setAlignment(Qt.AlignmentFlag.AlignLeft)

        # Combobox per la selezione della finestra o del monitor
        self.screenSelectionComboBox = CustomComboBox()
        self.screenSelectionComboBox.popupOpened.connect(self.updateWindowList)

        self.audioDeviceComboBox = QComboBox()
        audio_devices = self.print_audio_devices()
        self.audioDeviceComboBox.addItems(audio_devices)

        titles = [win.title for win in gw.getAllWindows() if win.title.strip()] + \
                 [f"Schermo intero {i + 1} - {w.width}x{w.height}" for i, w in enumerate(get_monitors())]
        self.screenSelectionComboBox.addItems(titles)

        # Campo di input per la scelta del nome del file e posizione
        self.filePathLineEdit = QLineEdit()
        self.filePathLineEdit.setPlaceholderText("Inserisci il percorso del file di destinazione")
        browseButton = QPushButton("Sfoglia")
        browseButton.clicked.connect(self.browseFileLocation)

        # Pulsanti per la registrazione
        self.startRecordingButton = QPushButton("")
        self.startRecordingButton.setIcon(QIcon("res/rec.png"))
        self.stopRecordingButton = QPushButton("")
        self.stopRecordingButton.setIcon(QIcon("res/stop.png"))
        buttonLayout = QHBoxLayout()
        buttonLayout.addWidget(self.startRecordingButton)
        buttonLayout.addWidget(self.stopRecordingButton)

        # Collegamento dei pulsanti ai metodi
        self.startRecordingButton.clicked.connect(self.startScreenRecording)
        self.stopRecordingButton.clicked.connect(self.stopScreenRecording)

        # Layout principale per il gruppo di controllo di registrazione
        recordingGroup = QGroupBox("Opzioni di Registrazione")
        recordingLayout = QVBoxLayout(recordingGroup)
        recordingLayout.addWidget(self.recordingStatusLabel)
        recordingLayout.addWidget(self.timecodeLabel)
        recordingLayout.addWidget(QLabel("Seleziona finestra o schermo:"))
        recordingLayout.addWidget(self.screenSelectionComboBox)
        recordingLayout.addWidget(QLabel("Seleziona input audio:"))
        recordingLayout.addWidget(self.audioDeviceComboBox)
        recordingLayout.addWidget(QLabel("File di destinazione:"))
        recordingLayout.addWidget(self.filePathLineEdit)
        recordingLayout.addWidget(browseButton)
        recordingLayout.addLayout(buttonLayout)

        # Widget per contenere il layout principale
        widget = QWidget()
        widget.setLayout(QVBoxLayout())
        widget.layout().addWidget(recordingGroup)
        dock.addWidget(widget)

        self.setDefaultAudioDevice()  # Imposta lo Stereo Mix come default subito dopo aver popolato la lista

        self.updateWindowList()
        return dock

    def updateWindowList(self):
        """Aggiorna la lista delle finestre e degli schermi disponibili, dando priorità agli schermi interi."""
        self.screenSelectionComboBox.clear()
        windows = [win.title for win in gw.getAllWindows() if win.title.strip()]

        # Ottieni i dettagli dei monitor e formatta il titolo per l'inserimento nella combo box
        monitors = [f"Schermo intero {i + 1} - {m.width}x{m.height}" for i, m in enumerate(get_monitors())]

        # Aggiungi prima i monitor alla lista della combo box
        self.screenSelectionComboBox.addItems(monitors + windows)

    def setDefaultAudioDevice(self):
        """Imposta 'Stereo Mix' come dispositivo predefinito se disponibile."""
        index = self.audioDeviceComboBox.findText("Stereo Mix")
        if index != -1:
            self.audioDeviceComboBox.setCurrentIndex(index)

    def browseFileLocation(self):
        """Apre un dialogo di selezione file per scegliere il percorso di salvataggio del video."""
        fileName, _ = QFileDialog.getSaveFileName(self, "Salva Video", "", "Video Files (*.avi)")
        if fileName:
            self.filePathLineEdit.setText(fileName)

        # Metodi per iniziare e fermare la registrazione

    def print_audio_devices(self):
        devices = sd.query_devices()
        available_audio_devices = []
        added_devices = set()  # Set per tenere traccia dei dispositivi già aggiunti

        for device in devices:
            # Verifica se il dispositivo è un microfono o lo Stereo Mix e non è già stato aggiunto
            if ('microphone' in device['name'].lower() or 'stereo mix' in device['name'].lower()) and device[
                'name'] not in added_devices:
                device_info = f"{device['name']} (Max canali: {device['max_input_channels']})"
                available_audio_devices.append(device_info)
                added_devices.add(device['name'])  # Aggiungi il nome del dispositivo al set

        return available_audio_devices

    def applyBackgroundAudioToVideo(self):
        video_path = self.videoPathLineEdit  # Percorso del video attualmente caricato
        background_audio_path = self.backgroundAudioPathLineEdit.text()  # Percorso dell'audio di sottofondo scelto
        background_volume = self.volumeSlider.value() / 100.0  # Volume dell'audio di sottofondo

        if not video_path or not os.path.exists(video_path):
            QMessageBox.warning(self, "Errore", "Carica un video prima di applicare l'audio di sottofondo.")
            return

        if not background_audio_path or not os.path.exists(background_audio_path):
            QMessageBox.warning(self, "Errore", "Carica un audio di sottofondo prima di applicarlo.")
            return

        try:
            # Carica video e audio di sottofondo
            video_clip = VideoFileClip(video_path)
            background_audio_clip = AudioFileClip(background_audio_path).volumex(background_volume)

            # Combina l'audio di sottofondo con l'audio originale del video, se presente
            if video_clip.audio:
                combined_audio = CompositeAudioClip([video_clip.audio, background_audio_clip])
            else:
                combined_audio = background_audio_clip

            # Imposta l'audio combinato nel video e salva il nuovo file
            final_clip = video_clip.set_audio(combined_audio)
            output_path = tempfile.mktemp(suffix='.mp4')
            final_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')

            QMessageBox.information(self, "Successo",
                                    f"Il video con audio di sottofondo è stato salvato in {output_path}")
            self.loadVideoOutput(output_path)  # Carica il video aggiornato nell'interfaccia
        except Exception as e:
            QMessageBox.critical(self, "Errore durante l'applicazione dell'audio di sottofondo", str(e))

    def applyAudioWithPauses(self):
        video_path = self.videoPathLineEdit  # Assicurati che questo sia il percorso del video attualmente caricato
        audio_path = self.audioPathLineEdit.text()
        pause_before = float(self.pauseBeforeLineEdit.text() or 0)
        pause_after = float(self.pauseAfterLineEdit.text() or 0)

        if not audio_path:
            # Estrai l'audio dal video corrente se non è stato fornito un percorso audio
            video_clip = VideoFileClip(video_path)
            audio_path = tempfile.mktemp(suffix='.mp3')  # Crea un percorso temporaneo per l'audio
            video_clip.audio.write_audiofile(audio_path)  # Salva l'audio estratto

        # Aggiungi pause e sostituisci l'audio nel video
        self.replaceAudioInVideo(video_path, audio_path, pause_before, pause_after)

    def updateTimecodeRec(self):
        if self.recordingTime is not None:
            self.recordingTime = self.recordingTime.addSecs(1)
            self.timecodeLabel.setText(self.recordingTime.toString("hh:mm:ss"))

    def startScreenRecording(self):
        selected_title = self.screenSelectionComboBox.currentText()
        selected_audio = self.audioDeviceComboBox.currentText()
        video_file_path = self.filePathLineEdit.text()

        if selected_title and selected_audio and video_file_path:
            # Utilizza regex per estrarre correttamente il nome del dispositivo
            match = re.match(r"^(.*?) \(Max canali:", selected_audio)
            if match:
                selected_audio_name = match.group(1)
                all_devices = sd.query_devices()
                audio_input_index = None
                for index, device in enumerate(all_devices):
                    if device['name'].startswith(selected_audio_name) and device['max_input_channels'] > 0:
                        audio_input_index = index
                        break

            if audio_input_index is not None:
                max_channels = all_devices[audio_input_index]['max_input_channels']
                audio_channels = min(2, max_channels)  # Usa un numero sicuro di canali

                monitors = get_monitors()
                if "Schermo intero" in selected_title:
                    index = int(selected_title.split()[2]) - 1
                    region = (monitors[index].x, monitors[index].y, monitors[index].width, monitors[index].height)
                else:
                    window = gw.getWindowsWithTitle(selected_title)[0]
                    window.moveTo(0, 0)
                    window.activate()
                    region = (window.left, window.top, window.width, window.height)

                # Genera un timestamp per creare un nome unico per ogni file
                timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
                base_filename = os.path.splitext(video_file_path)[0]
                extension = os.path.splitext(video_file_path)[1]
                video_file_path_with_timestamp = f"{base_filename}_{timestamp}{extension}"
                audioFileName = f"{base_filename}_{timestamp}.wav"

                fourcc = cv2.VideoWriter_fourcc(*'XVID')
                self.video_writer = cv2.VideoWriter(video_file_path_with_timestamp, fourcc, 25.0,
                                                    (region[2], region[3]))

                self.recorder_thread = ScreenRecorder(self.video_writer, audioFileName, region=region,
                                                      audio_input=audio_input_index, audio_channels=audio_channels)
                self.recorder_thread.start()

                self.recordingStatusLabel.setText("Stato: Registrazione in corso")
                self.recordingTime = QTime(0, 0, 0)
                self.timer.start(1000)  # Avvia il timer che aggiorna ogni secondo

                self.current_video_path = video_file_path_with_timestamp
                self.current_audio_path = audioFileName
                self.recordingStatusLabel.setText(f'Stato: Registrazione iniziata di {selected_title}')

            else:
                QMessageBox.warning(self, "Errore", "Dispositivo audio non trovato.")
        else:
            QMessageBox.warning(self, "Errore",
                                "Assicurati di selezionare una finestra/schermo, un dispositivo audio e un percorso di salvataggio valido.")

    def stopScreenRecording(self):
        # Stop the recording process
        if self.recorder_thread is not None:
            self.recorder_thread.stop()

        # Close the video writer and release resources
        if self.video_writer is not None:
            self.video_writer.release()
            self.video_writer = None

        # Example of handling the paths set earlier
        if self.current_video_path and self.current_audio_path:
            # Process to merge audio and video or finalize the recording
            self.mergeAudioVideo(self.current_video_path, self.current_audio_path)

        self.timer.stop()
        self.recordingStatusLabel.setText("Stato: Registrazione Terminata.")

    def mergeAudioVideo(self, video_path, audio_path):
        try:
            video_clip = VideoFileClip(video_path)
            audio_clip = AudioFileClip(audio_path)
            final_clip = video_clip.set_audio(audio_clip)
            output_path = video_path.replace('.avi', '_final.mp4')
            final_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')
            QMessageBox.information(self, "Successo", f"Video finale salvato in: {output_path}")
            self.loadVideoOutput(output_path)  # Carica il video finale

        except Exception as e:
            QMessageBox.critical(self, "Errore", f"Errore durante l'unione di audio e video: {e}")

    def saveText(self):
        # Apri il dialogo di salvataggio file e ottieni il percorso del file e il filtro selezionato dall'utente
        path, _ = QFileDialog.getSaveFileName(self, "Salva file", "", "Text files (*.txt);;All files (*.*)")

        # Controlla se l'utente ha effettivamente scelto un file
        if path:
            # Ottieni il testo dal QTextEdit
            text_to_save = self.transcriptionTextArea.toPlainText()

            # Prova a salvare il testo nel file scelto
            try:
                with open(path, 'w') as file:
                    file.write(text_to_save)
                print("File salvato correttamente!")
            except Exception as e:
                print("Errore durante il salvataggio del file:", e)

    def createDownloadDock(self):
        """Crea e restituisce il dock per il download di video."""
        dock = Dock("Download Video")

        # GroupBox per organizzare visivamente le opzioni di download
        downloadGroup = QGroupBox("Opzioni di Download Video")
        downloadLayout = QVBoxLayout()

        # Widget per inserire l'URL del video di YouTube
        url_label = QLabel("Inserisci l'URL di YouTube:")
        url_edit = QLineEdit()
        download_btn = QPushButton("Download Video")
        download_btn.clicked.connect(lambda: self.handleDownload(url_edit.text()))

        # Aggiunta dei controlli al layout della GroupBox
        downloadLayout.addWidget(url_label)
        downloadLayout.addWidget(url_edit)
        downloadLayout.addWidget(download_btn)

        # Imposta il layout del GroupBox
        downloadGroup.setLayout(downloadLayout)

        # Widget principale per il dock
        widget = QWidget()
        widget.setLayout(QVBoxLayout())
        widget.layout().addWidget(downloadGroup)
        dock.addWidget(widget)

        return dock

    def handleDownload(self, url):
        if url:
            self.downloadThread = DownloadThread(url)
            self.downloadThread.finishedAudio.connect(self.onDownloadFinished)
            self.downloadThread.error.connect(self.onError)
            self.downloadThread.progress.connect(self.updateDownloadProgress)  # Connect to the new progress signal
            self.downloadThread.start()
            self.showDownloadProgress()

    def showDownloadProgress(self):
        self.progressDialog = QProgressDialog("Downloading video...", "Abort", 0, 100, self)
        self.progressDialog.setWindowTitle("Download Progress")
        self.progressDialog.setWindowModality(Qt.WindowModality.WindowModal)
        self.progressDialog.canceled.connect(self.downloadThread.terminate)  # Connect cancel action
        self.progressDialog.show()

    def updateDownloadProgress(self, progress):
        if not self.progressDialog.wasCanceled():
            self.progressDialog.setValue(progress)

    def onDownloadFinished(self, audio_path, video_title, video_language):
        self.progressDialog.close()
        QMessageBox.information(self, "Download Complete", f"audio saved to {audio_path}.")
        self.video_download_language = video_language
        print (video_language)
        if self.isAudioOnly(audio_path):
            self.loadVideo(audio_path, video_title)

    def onError(self, error_message):
        self.progressDialog.close()
        QMessageBox.critical(self, "Download Error", error_message)

    def isAudioOnly(self, file_path):
        """Check if the file is likely audio-only based on the extension."""
        audio_extensions = {'.mp3', '.wav', '.aac', '.m4a', '.flac', '.ogg'}
        ext = os.path.splitext(file_path)[1].lower()
        return ext in audio_extensions

    def sourceSetter(self, url):
        self.player.setSource(QUrl.fromLocalFile(url))
        #self.player.play()

    def sourceSetterOutput(self, url):
        self.playerOutput.setSource(QUrl.fromLocalFile(url))
        self.playerOutput.play()

    def loadVideo(self, video_path, video_title = 'Video Track'):
        """Load and play video or audio, updating UI based on file type."""
        # Scarica il video corrente prima di caricarne uno nuovo
        self.player.stop()

        if self.player.playbackState() == QMediaPlayer.PlaybackState.StoppedState:
            QTimer.singleShot(1, lambda: self.sourceSetter(video_path))

        self.videoPathLineEdit = video_path

        if self.isAudioOnly(video_path):
            self.fileNameLabel.setText(f"{video_title} - Traccia solo audio")  # Display special message for audio files
        else:
            self.fileNameLabel.setText(os.path.basename(video_path))  # Update label with file name

        self.updateRecentFiles(video_path)  # Update recent files list

    def loadVideoOutput(self, video_path):

        self.playerOutput.stop()

        if self.playerOutput.playbackState() == QMediaPlayer.PlaybackState.StoppedState:
            QTimer.singleShot(1, lambda: self.sourceSetterOutput(video_path))

        self.fileNameLabelOutput.setText(os.path.basename(video_path))  # Aggiorna il nome del file sulla label
        self.videoPathLineOutputEdit = video_path
        print(f"Loaded video output: {video_path}")


    def updateTimeCode(self, position):
        # Calcola ore, minuti e secondi dalla posizione, che è in millisecondi
        hours, remainder = divmod(position // 1000, 3600)
        minutes, seconds = divmod(remainder, 60)
        # Aggiorna l'etichetta con il nuovo time code
        self.currentTimeLabel.setText(f'{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}')

    def updateDuration(self, duration):
        # Calcola ore, minuti e secondi dalla durata, che è in millisecondi
        hours, remainder = divmod(duration // 1000, 3600)
        minutes, seconds = divmod(remainder, 60)
        # Aggiorna l'etichetta con la durata totale
        self.totalTimeLabel.setText(f' / {int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}')

    def setupMenuBar(self):
        menuBar = self.menuBar()
        fileMenu = menuBar.addMenu('&File')

        openAction = QAction('&Open Video/Audio', self)
        openAction.setShortcut('Ctrl+O')
        openAction.setStatusTip('Open video')
        openAction.triggered.connect(self.browseVideo)


        openActionOutput = QAction('&Open as Output Video', self)
        openAction.setShortcut('Ctrl+I')
        openActionOutput.setStatusTip('Open Video Output')
        openActionOutput.triggered.connect(self.browseVideoOutput)

        fileMenu.addAction(openAction)
        fileMenu.addAction(openActionOutput)

        # Azione per generare la presentazione PowerPoint
        generatePresentationAction = QAction('Genera &Presentazione', self)
        generatePresentationAction.setShortcut('Ctrl+P')
        generatePresentationAction.setStatusTip('Genera una presentazione PowerPoint')
        generatePresentationAction.triggered.connect(self.generaPresentationConTestoAttuale)
        fileMenu.addAction(generatePresentationAction)

        exitAction = QAction('&Exit', self)
        exitAction.setShortcut('Ctrl+Q')
        exitAction.setStatusTip('Exit application')
        exitAction.triggered.connect(self.close)
        fileMenu.addAction(exitAction)
        fileMenu.addSeparator()

        self.recentMenu = fileMenu.addMenu("Recenti")  # Aggiunge il menu dei file recenti
        self.updateRecentFilesMenu()
        # Creazione del menu View per la gestione della visibilità dei docks
        viewMenu = menuBar.addMenu('&View')
        self.setupViewMenuActions(viewMenu)
        # Creazione del menu About
        aboutMenu = menuBar.addMenu('&About')
        # Aggiunta di azioni al menu About
        aboutAction = QAction('&About', self)
        aboutAction.setStatusTip('About the application')
        aboutAction.triggered.connect(self.about)
        aboutMenu.addAction(aboutAction)
        self.updateViewMenu()

    def setupViewMenuActions(self, viewMenu):
        # Azione per il Video Player Dock
        self.actionToggleVideoPlayerDock = QAction('Mostra/Nascondi Video Player Source', self, checkable=True)
        self.actionToggleVideoPlayerDock.setChecked(self.videoPlayerDock.isVisible())
        self.actionToggleVideoPlayerDock.triggered.connect(
            lambda: self.toggleDockVisibilityAndUpdateMenu(self.videoPlayerDock,
                                                           self.actionToggleVideoPlayerDock.isChecked()))

        # Azioni simili per gli altri docks...
        self.actionToggleVideoPlayerDockOutput = self.createToggleAction(self.videoPlayerOutputDock,
                                                                         'Mostra/Nascondi Video Player Output')
        self.actionToggleTranscriptionDock = self.createToggleAction(self.transcriptionDock,
                                                                     'Mostra/Nascondi Trascrizione')
        self.actionToggleEditingDock = self.createToggleAction(self.editingDock, 'Mostra/Nascondi Editing')
        self.actionToggleDownloadDock = self.createToggleAction(self.downloadDock, 'Mostra/Nascondi Download')
        self.actionToggleRecordingDock = self.createToggleAction(self.recordingDock, 'Mostra/Nascondi Registrazione')
        self.actionToggleAudioDock = self.createToggleAction(self.audioDock, 'Mostra/Nascondi Gestione Audio')
        self.actionToggleVideoMergeDock = self.createToggleAction(self.videoMergeDock, 'Mostra/Nascondi Unisci Video')

        # Aggiungi tutte le azioni al menu 'View'
        viewMenu.addAction(self.actionToggleVideoPlayerDock)
        viewMenu.addAction(self.actionToggleVideoPlayerDockOutput)
        viewMenu.addAction(self.actionToggleTranscriptionDock)
        viewMenu.addAction(self.actionToggleEditingDock)
        viewMenu.addAction(self.actionToggleDownloadDock)
        viewMenu.addAction(self.actionToggleRecordingDock)
        viewMenu.addAction(self.actionToggleAudioDock)
        viewMenu.addAction(self.actionToggleVideoMergeDock)


        # Aggiungi azioni per mostrare/nascondere tutti i docks
        showAllDocksAction = QAction('Mostra tutti i Docks', self)
        hideAllDocksAction = QAction('Nascondi tutti i Docks', self)

        showAllDocksAction.triggered.connect(self.showAllDocks)
        hideAllDocksAction.triggered.connect(self.hideAllDocks)

        viewMenu.addSeparator()  # Aggiunge un separatore per chiarezza
        viewMenu.addAction(showAllDocksAction)
        viewMenu.addAction(hideAllDocksAction)


        # Azione per salvare il layout dei docks
        saveLayoutAction = QAction('Salva Layout dei Docks', self)
        saveLayoutAction.triggered.connect(self.saveDockLayout)
        viewMenu.addSeparator()  # Aggiunge un separatore per chiarezza
        viewMenu.addAction(saveLayoutAction)

        # Aggiorna lo stato iniziale del menu
        self.updateViewMenu()

    def saveDockLayout(self):
        # Assumendo che DockSettingsManager abbia un metodo per salvare le impostazioni
        if hasattr(self, 'dockSettingsManager'):
            self.dockSettingsManager.save_settings()
            QMessageBox.information(self, "Layout Salvato", "Il layout dei docks è stato salvato correttamente.")
        else:
            QMessageBox.warning(self, "Errore", "Gestore delle impostazioni dei dock non trovato.")

    def showAllDocks(self):
        # Imposta tutti i docks visibili
        self.videoPlayerDock.setVisible(True)
        self.videoPlayerOutputDock.setVisible(True)
        self.audioDock.setVisible(True)
        self.transcriptionDock.setVisible(True)
        self.editingDock.setVisible(True)
        self.downloadDock.setVisible(True)
        self.recordingDock.setVisible(True)
        self.videoMergeDock.setVisible(True)
        self.updateViewMenu()  # Aggiorna lo stato dei menu

    def hideAllDocks(self):
        # Nasconde tutti i docks
        self.videoPlayerDock.setVisible(False)
        self.videoPlayerOutputDock.setVisible(False)
        self.audioDock.setVisible(False)
        self.transcriptionDock.setVisible(False)
        self.editingDock.setVisible(False)
        self.downloadDock.setVisible(False)
        self.recordingDock.setVisible(False)
        self.videoMergeDock.setVisible(False)
        self.updateViewMenu()  # Aggiorna lo stato dei menu
    def createToggleAction(self, dock, menuText):
        action = QAction(menuText, self, checkable=True)
        action.setChecked(dock.isVisible())
        action.triggered.connect(lambda checked: self.toggleDockVisibilityAndUpdateMenu(dock, checked))
        return action

    def toggleDockVisibilityAndUpdateMenu(self, dock, visible):
        dock.setVisible(visible)
        self.updateViewMenu()

    def updateViewMenu(self):
        # Aggiorna lo stato dei menu checkable basato sulla visibilità dei dock
        self.actionToggleVideoPlayerDock.setChecked(self.videoPlayerDock.isVisible())
        self.actionToggleVideoPlayerDockOutput.setChecked(self.videoPlayerOutputDock.isVisible())
        self.actionToggleAudioDock.setChecked(self.audioDock.isVisible())
        self.actionToggleTranscriptionDock.setChecked(self.transcriptionDock.isVisible())
        self.actionToggleEditingDock.setChecked(self.editingDock.isVisible())
        self.actionToggleDownloadDock.setChecked(self.downloadDock.isVisible())
        self.actionToggleRecordingDock.setChecked(self.recordingDock.isVisible())
        self.actionToggleVideoMergeDock.setChecked(self.videoMergeDock.isVisible())

    def about(self):
        QMessageBox.about(self, "TGeniusAI",
                          """<b>Thema Genius</b> version 1.0<br>
                          AI.<br>
                          <br>
                          Autore: FFA <br>""")

    def convertVideoToAudio(self, video_file, audio_format='wav'):
        """Estrae la traccia audio dal video e la converte in formato WAV."""
        with VideoFileClip(video_file) as video_clip:
            # Crea un file temporaneo nel sistema in modo sicuro
            temp_audio = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{audio_format}')
            # Esegui l'estrazione dell'audio e salvalo nel file temporaneo
            video_clip.audio.write_audiofile(temp_audio.name, codec='pcm_s16le')  # codec per wav
            # Chiudi il file temporaneo per evitare lock sul file
            temp_audio.close()
            return temp_audio.name

    def splitAudio(self, audio_file, length=60000):
        """Divide l'audio in blocchi di una durata specifica (in millisecondi)."""
        audio = AudioSegment.from_file(audio_file)
        chunks = [(audio[i:i + length], i) for i in range(0, len(audio), length)]  # Includere il timestamp di inizio
        return chunks

    def transcribeAudioChunk(self, audio_chunk, start_time):
        def get_locale_from_language(language_code):
            """Converte un codice di lingua ISO 639-1 in un locale più specifico."""
            try:
                language = pycountry.languages.get(alpha_2=language_code)
                # Mappatura semplificata: mappa 'en' a 'en-US', ecc.
                return {
                    'en': 'en-US',
                    'es': 'es-ES',
                    'fr': 'fr-FR',
                    'it': 'it-IT',
                    'de': 'de-DE'
                }.get(language.part1, f"{language.part1}-{language.part1.upper()}")
            except Exception:
                return language_code  # Ritorna il codice originale se la mappatura fallisce

        def try_remove_chunk_file(chunk_file):
            """Rimuove il file audio temporaneo con tentativi multipli."""
            max_attempts = 3
            for attempt in range(max_attempts):
                try:
                    os.remove(chunk_file)
                    break
                except PermissionError:
                    if attempt < max_attempts - 1:
                        time.sleep(0.5)  # Aspetta un po' prima di riprovare

        recognizer = sr.Recognizer()
        unique_id = uuid.uuid4()
        chunk_file = f"temp_chunk_{unique_id}.wav"
        audio_chunk.export(chunk_file, format="wav")

        try:
            with sr.AudioFile(chunk_file) as source:
                audio_data = recognizer.record(source)
                language_video = self.video_download_language if self.video_download_language else self.languageComboBox.currentData()
                language_video = get_locale_from_language(
                    language_video)  # Usa la funzione per ottenere il locale corretto
                text = recognizer.recognize_google(audio_data, language=language_video)
            return text, start_time, language_video
        except sr.UnknownValueError:
            return "[Incomprensibile]", start_time, language_video
        except sr.RequestError as e:
            return f"[Errore: {e}]", start_time, language_video
        finally:
            try_remove_chunk_file(chunk_file)

    def onLanguageChange(self):
        # Ottieni il codice lingua della selezione corrente
        language_code = self.languageComboBox.currentData()
        print(f"Lingua selezionata cambiata a: {language_code}")
        # Puoi aggiungere qui ulteriori operazioni basate sul cambio di lingua

    def updateLanguageComboBox(self, language_code, language_name):
        # Verifica se la lingua è già presente nella combo box
        index = self.languageComboBox.findData(language_code)
        if index == -1:  # Lingua non presente, aggiungila
            self.languageComboBox.addItem(language_name, language_code)
            index = self.languageComboBox.count() - 1
        self.languageComboBox.setCurrentIndex(index)
    def handleTextChange(self):
        text = self.transcriptionTextArea.toPlainText()
        if text.strip():
            try:
                detected_language_code = detect(text)
                language = pycountry.languages.get(alpha_2=detected_language_code)
                if language:
                    detected_language = language.name
                    self.updateLanguageComboBox(detected_language_code, detected_language)
                    self.updateTranscriptionLanguageDisplay(detected_language)
                else:
                    # If the detected language is not supported by pycountry, display a default message.
                    self.updateTranscriptionLanguageDisplay("Lingua non supportata")
            except LangDetectException:
                self.updateTranscriptionLanguageDisplay("Non rilevabile")
        else:
            self.languageComboBox.setCurrentIndex(-1)  # Resetta la selezione se non c'è testo
            self.updateTranscriptionLanguageDisplay("")

    def updateTranscriptionLanguageDisplay(self, language):
        """
        Aggiorna il dock della trascrizione con la lingua attuale.
        """
        self.transcriptionLanguageLabel.setText(f"Lingua rilevata: {language}")
    def transcribeVideo(self):
        if not self.videoPathLineEdit:
            QMessageBox.warning(self, "Attenzione", "Nessun video selezionato.")
            return

        self.progressDialog = QProgressDialog("Trascrizione in corso...", "Annulla", 0, 100, self)
        self.progressDialog.setWindowTitle("Progresso Trascrizione")
        self.progressDialog.setWindowModality(Qt.WindowModality.WindowModal)
        self.progressDialog.show()

        self.thread = TranscriptionThread(self.videoPathLineEdit, self)
        self.thread.update_progress.connect(self.updateProgressDialog(self.progressDialog))
        self.thread.transcription_complete.connect(self.completeTranscription(self.progressDialog))
        self.thread.error_occurred.connect(self.handleErrors(self.progressDialog))
        self.thread.start()

    def updateProgressDialog(self, progress_dialog):
        def update(value, label):
            if not progress_dialog.wasCanceled():
                progress_dialog.setValue(value)
                progress_dialog.setLabelText(label)

        return update

    def completeTranscription(self, progress_dialog):
        """Returns a closure that handles transcription completion."""

        def complete(text, temp_files):
            if not progress_dialog.wasCanceled():
                self.transcriptionTextArea.setText(text)
                progress_dialog.setValue(100)
                progress_dialog.close()
                self.cleanupFiles(temp_files)  # Immediately cleanup after transcription

        return complete

    def cleanupFiles(self, file_paths):
        """Safely removes temporary files used during transcription."""
        for path in file_paths:
            self.removeFileSafe(path)

    def removeFileSafe(self, file_path, attempts=5, delay=0.5):
        """Attempt to safely remove a file with retries and delays."""
        for _ in range(attempts):
            try:
                os.remove(file_path)
                print(f"File {file_path} successfully removed.")
                break
            except PermissionError:
                print(f"Warning: File {file_path} is currently in use. Retrying...")
                time.sleep(delay)
            except FileNotFoundError:
                print(f"The file {file_path} does not exist or has already been removed.")
                break
            except Exception as e:
                print(f"Unexpected error while removing {file_path}: {e}")
    def handleErrors(self, progress_dialog):
        def error(message):
            QMessageBox.critical(self, "Errore nella Trascrizione",
                                 f"Errore durante la trascrizione del video: {message}")
            progress_dialog.cancel()

        return error



    def impostaFont(shape, size_pt, text):
        """
        Imposta il font per una shape data con la grandezza specificata.
        """
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        font = run.font
        font.size = Pt(size_pt)  # Imposta la grandezza del font



    def generateAudioWithElevenLabs(self):
        def convert_numbers_to_words(text):
            # Modifica il testo in modo che ogni numero seguito direttamente da un punto sia seguito da uno spazio
            text = re.sub(r'(\d+)\.', r'\1 .', text)

            new_text = []
            for word in text.split():
                if word.isdigit():
                    new_word = num2words(word, lang='it')
                    new_text.append(new_word)
                else:
                    new_text.append(word)
            return ' '.join(new_text)

        transcriptionText = self.transcriptionTextArea.toPlainText()
        if not transcriptionText.strip():
            QMessageBox.warning(self, "Attenzione", "Inserisci una trascrizione prima di generare l'audio.")
            return
        transcriptionText = convert_numbers_to_words(transcriptionText)
        # Dati per identificare la voce e il modello
        voice_id = self.voiceSelectionComboBox.currentData()
        model_id = "eleven_multilingual_v1"  # Assumi il modello, personalizza come necessario

        voice_settings = {
            'stability': self.stabilitySlider.value() / 100.0,
            'similarity_boost': self.similaritySlider.value() / 100.0,
            'style': self.styleSlider.value() / 10.0,
            'use_speaker_boost': self.speakerBoostCheckBox.isChecked()
        }

        self.progressDialog = QProgressDialog("Generazione audio in corso...", "Annulla", 0, 100, self)
        self.progressDialog.setWindowTitle("Progresso Generazione Audio")
        self.progressDialog.setWindowModality(Qt.WindowModality.WindowModal)

        # Crea il thread con i nuovi parametri
        self.audio_thread = AudioGenerationThread(transcriptionText, voice_id, model_id, voice_settings,
                                                  "ef38b436326ec387ecb1a570a8641b84", self)
        self.audio_thread.progress.connect(self.progressDialog.setValue)
        self.audio_thread.completed.connect(self.onAudioGenerationCompleted)
        self.audio_thread.error.connect(self.onError)
        self.audio_thread.start()

        # Prepara il dialogo di progresso
        self.progressDialog.canceled.connect(self.audio_thread.terminate)
        self.progressDialog.show()

    def onAudioGenerationCompleted(self, audio_path):
        QMessageBox.information(self, "Generazione Completata", f"L'audio è stato salvato in: {audio_path}")

    def addPauseAndMerge(self, original_audio_path, pause_before, pause_after):
        # Carica il file audio originale usando pydub
        original_audio = AudioSegment.from_file(original_audio_path)

        # Crea clip audio di silenzio per le pause
        pause_before_clip = AudioSegment.silent(duration=int(pause_before * 1000))  # durata in millisecondi
        pause_after_clip = AudioSegment.silent(duration=int(pause_after * 1000))

        # Concatena le pause con l'audio originale
        final_audio = pause_before_clip + original_audio + pause_after_clip

        # Salva il nuovo file audio
        new_audio_path = tempfile.mktemp(suffix='.mp3')
        final_audio.export(new_audio_path, format='mp3')

        return new_audio_path

    def onAudioGenerationCompleted(self, audio_path):

        # Gestione della durata delle pause come prima
        durata_pausa_iniziale = float(self.pauseBeforeLineEdit.text() or 0)
        durata_pausa_finale = float(self.pauseAfterLineEdit.text() or 0)

        if durata_pausa_iniziale > 0 or durata_pausa_finale > 0:
            audio_path = self.addPauseAndMerge(audio_path, durata_pausa_iniziale, durata_pausa_finale)

        # Continua con le operazioni di combinazione audio/video come prima
        base_name = os.path.splitext(os.path.basename(self.videoPathLineEdit))[0]
        timestamp = time.strftime('%Y%m%d%H%M%S', time.localtime())
        output_path = os.path.join(os.path.dirname(self.videoPathLineEdit), f"{base_name}_GeniusAI_{timestamp}.mp4")
        self.adattaVelocitaVideoAAudio(self.videoPathLineEdit, audio_path, output_path)

        try:
            if os.path.exists(audio_path):
                os.remove(audio_path)
        except Exception as e:
            QMessageBox.critical(self, "Errore", f"Non è stato possibile eliminare il file audio temporaneo: {e}")

        QMessageBox.information(self, "Completato", "Processo completato con successo!")

        self.loadVideoOutput(output_path)

    def onError(self, error_message):
        QMessageBox.critical(self, "Errore", "Errore durante la generazione dell'audio: " + error_message)

    def mergeAudioTracks(self, lista_percorsi_audio, output_path):
        # Assicurati che la lista non sia vuota
        if not lista_percorsi_audio:
            raise ValueError("La lista dei percorsi audio è vuota")

        # Carica la prima traccia audio
        traccia_unita = AudioSegment.from_file(lista_percorsi_audio[0])

        # Unisci le tracce rimanenti
        for percorso in lista_percorsi_audio[1:]:
            traccia_attuale = AudioSegment.from_file(percorso)
            traccia_unita += traccia_attuale

        # Esporta la traccia audio unita
        traccia_unita.export(output_path, format="mp3")

    def browseAudio(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Seleziona Audio", "", "Audio Files (*.mp3 *.wav)")
        if fileName:
            self.audioPathLineEdit.setText(fileName)  # Aggiorna il campo di testo con il percorso del file

    def extractAudioFromVideo(self, video_path):
        # Estrai l'audio dal video e salvalo temporaneamente
        temp_audio_path = tempfile.mktemp(suffix='.mp3')
        video_clip = VideoFileClip(video_path)
        video_clip.audio.write_audiofile(temp_audio_path)
        return temp_audio_path

    def applyNewAudioToVideo(self, video_path, audio_path):
        # Carica il video e l'audio modificato
        video_clip = VideoFileClip(video_path)
        audio_clip = AudioFileClip(audio_path)
        final_clip = video_clip.set_audio(audio_clip)
        # Salvataggio del video finale
        output_video_path = video_path.replace('.mp4', '_new_audio.mp4')
        final_clip.write_videofile(output_video_path, codec='libx264', audio_codec='aac')

        # Aggiorna l'interfaccia utente per riflettere il cambio
        self.loadVideoOutput(output_video_path)

    def replaceAudioInVideo(self, video_path, audio_path, pause_before, pause_after):
        if not audio_path:
            audio_path = self.extractAudioFromVideo(video_path)

        # Aggiungi pause all'audio
        new_audio_path = self.addPauseAndMerge(audio_path, pause_before, pause_after)

        video_clip = None
        new_audio_clip = None
        adjusted_video_clip = None
        try:
            # Sostituisci l'audio nel video e adatta la velocità del video all'audio modificato
            video_clip = VideoFileClip(video_path)
            new_audio_clip = AudioFileClip(new_audio_path)

            # Adatta la velocità del video alla durata dell'audio modificato
            output_video_path = video_path.replace('.mp4', '_adjusted_audio.mp4')
            self.adattaVelocitaVideoAAudio(video_path, new_audio_path, output_video_path)

            # Imposta il nuovo audio nel video velocizzato
            adjusted_video_clip = VideoFileClip(output_video_path)
            final_clip = adjusted_video_clip.set_audio(new_audio_clip)
            final_output_video_path = output_video_path.replace('_adjusted_audio.mp4', '_final.mp4')
            final_clip.write_videofile(final_output_video_path, codec='libx264', audio_codec='aac')

            QMessageBox.information(self, "Successo",
                                    "Audio con pause applicato e velocità del video adattata con successo.")
        except Exception as e:
            QMessageBox.critical(self, "Errore", f"Errore durante la sostituzione dell'audio: {e}")
        finally:
            # Chiudi i clip per liberare le risorse
            if video_clip:
                video_clip.close()
            if new_audio_clip:
                new_audio_clip.close()
            if adjusted_video_clip:
                adjusted_video_clip.close()

            # Pulizia: rimuovi i file audio temporanei se necessario
            if os.path.exists(new_audio_path):
                os.remove(new_audio_path)
            if os.path.exists(output_video_path) and output_video_path != final_output_video_path:
                os.remove(output_video_path)

            # Aggiorna il video output per riflettere le modifiche
            self.loadVideoOutput(final_output_video_path)

    def cutVideo(self):
        media_path = self.videoPathLineEdit
        if not media_path:
            QMessageBox.warning(self, "Attenzione", "Per favore, seleziona un file prima di tagliarlo.")
            return

        if media_path.lower().endswith(('.mp4', '.mov', '.avi')):
            is_audio = False
        elif media_path.lower().endswith(('.mp3', '.wav', '.aac', '.ogg', '.flac')):
            is_audio = True
        else:
            QMessageBox.warning(self, "Errore", "Formato file non supportato.")
            return

        start_time = self.currentPosition / 1000.0  # Converti in secondi

        base_name = os.path.splitext(os.path.basename(media_path))[0]
        directory = os.path.dirname(media_path)
        ext = 'mp4' if not is_audio else 'mp3'
        output_path1 = os.path.join(directory, f"{base_name}_part1.{ext}")
        output_path2 = os.path.join(directory, f"{base_name}_part2.{ext}")

        self.progressDialog = QProgressDialog("Taglio del file in corso...", "Annulla", 0, 100, self)
        self.progressDialog.setWindowTitle("Progresso Taglio")
        self.progressDialog.setWindowModality(Qt.WindowModality.WindowModal)
        self.progressDialog.show()

        self.cutting_thread = VideoCuttingThread(media_path, start_time, output_path1, output_path2)
        self.cutting_thread.progress.connect(self.progressDialog.setValue)
        self.cutting_thread.completed.connect(self.onCutCompleted)
        self.cutting_thread.error.connect(self.onCutError)


        self.cutting_thread.start()

    def onCutCompleted(self, part1, part2):
        QMessageBox.information(self, "Successo", f"File tagliato e salvato in due parti: {part1} e {part2}.")
        self.progressDialog.close()

    def onCutError(self, error_message):
        QMessageBox.critical(self, "Errore", error_message)
        self.progressDialog.close()

    def positionChanged(self, position):
        self.videoSlider.setValue(position)
        self.currentPosition = position  # Aggiorna la posizione corrente
        self.updateTimeCode(position)

    # Slot per aggiornare il range massimo dello slider in base alla durata del video
    def durationChanged(self, duration):
        self.videoSlider.setRange(0, duration)
        self.updateDuration(duration)

    # Slot per aggiornare la posizione dello slider in base alla posizione corrente del video

    # Slot per cambiare la posizione del video quando lo slider viene mosso
    def setPosition(self, position):
        self.player.setPosition(position)
    def applyDarkMode(self):
        self.setStyleSheet("""
            QWidget {
                background-color: #2b2b2b;
                color: #dcdcdc;
            }
            QLineEdit {
                background-color: #333333;
                border: 1px solid #555555;
                border-radius: 2px;
                padding: 5px;
                color: #ffffff;
            }
            QPushButton {
                background-color: #555555;
                border: 1px solid #666666;
                border-radius: 2px;
                padding: 5px;
                color: #ffffff;
            }
            QPushButton:hover {
                background-color: #666666;
            }
            QPushButton:pressed {
                background-color: #777777;
            }
            QLabel {
                color: #cccccc;
            }
            QFileDialog {
                background-color: #444444;
            }
            QMessageBox {
                background-color: #444444;
            }
        """)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event):
        file_urls = [url.toLocalFile() for url in event.mimeData().urls() if url.isLocalFile()]
        if file_urls:

            self.videoPathLineEdit = file_urls[0]  # Aggiorna il percorso del video memorizzato
            self.loadVideo(self.videoPathLineEdit, os.path.basename(file_urls[0]))

    def browseVideo(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Seleziona Video", "", "Video/Audio Files (*avi *.mp4 *.mov *.mp3 *.wav *.aac *.ogg *.flac)")
        if fileName:
           self.loadVideo(fileName)

    def browseVideoOutput(self):
        fileName, _ = QFileDialog.getOpenFileName(self, "Seleziona Video", "", "Video/Audio Files (*.avi *.mp4 *.mov *.mp3 *.wav *.aac *.ogg *.flac)")
        if fileName:
           self.loadVideoOutput(fileName)

    def updateRecentFiles(self, newFile):
        if newFile not in self.recentFiles:
            self.recentFiles.insert(0, newFile)
            if len(self.recentFiles) > 5:  # Limita la lista ai 5 più recenti
                self.recentFiles.pop()
        self.updateRecentFilesMenu()

    def updateRecentFilesMenu(self):
        self.recentMenu.clear()  # Pulisce le voci precedenti
        for file in self.recentFiles:
            action = QAction(os.path.basename(file), self)
            action.triggered.connect(lambda checked, f=file: self.openRecentFile(f))
            self.recentMenu.addAction(action)

    def openRecentFile(self, filePath):
        self.videoPathLineEdit = filePath
        self.player.setSource(QUrl.fromLocalFile(filePath))
        self.fileNameLabel.setText(os.path.basename(filePath))
    def playVideo(self):
        self.player.play()

    def pauseVideo(self):
        self.player.pause()

    def adattaVelocitaVideoAAudio(self, video_path, new_audio_path, output_path):
        try:
            # Carica il nuovo file audio e calcola la sua durata
            new_audio = AudioFileClip(new_audio_path)
            durata_audio = new_audio.duration

            # Carica il video (senza audio) e calcola la sua durata
            video_clip = VideoFileClip(video_path)
            durata_video = video_clip.duration

            # Calcola il fattore di velocità necessario per far combaciare le durate
            fattore_velocita = durata_video / durata_audio

            # Applica il fattore di velocità al video
            video_modificato = video_clip.fx(vfx.speedx, fattore_velocita)

            # Imposta il nuovo audio sul video modificato
            final_video = video_modificato.set_audio(new_audio)

            # Scrivi il video finale mantenendo lo stesso frame rate del video originale
            original_frame_rate = video_clip.fps

            # Specifica del codec libx264 per il video e aac per l'audio
            final_video.write_videofile(output_path, codec="libx264", audio_codec="aac", fps=original_frame_rate)
            print(f'Video elaborato con successo.')
        except Exception as e:
            print(f"Errore durante l'adattamento della velocità del video: {e}")
    def stopVideo(self):
        self.player.stop()

    def unisci_video(self, lista_percorsi_video, percorso_file_output):
        try:
            clips = [VideoFileClip(video) for video in lista_percorsi_video]
            video_finale = concatenate_videoclips(clips)
            video_finale.write_videofile(percorso_file_output, codec="libx264", audio_codec="aac")
            QMessageBox.information(self, "Successo", f"Video unito creato con successo in {percorso_file_output}.")
        except Exception as e:
            QMessageBox.critical(self, "Errore", f"Errore durante l'unione dei video: {e}")

    def creaPresentazione(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Seleziona File di Testo", "", "Text Files (*.txt)")
        if file_path:
            try:
                # Qui chiamiamo la funzione `crea_presentation_da_file`, passando il percorso del file selezionato
                self.createPresentationFromFile(file_path)
                QMessageBox.information(self, "Successo", "Presentazione PowerPoint generata con successo.")
            except Exception as e:
                QMessageBox.critical(self, "Errore",
                                     f"Si è verificato un errore durante la generazione della presentazione: {e}")

    def createPresentationFromFile(self, file_path):
        # Tentativo di lettura del file con diverse codifiche
        encodings = ['utf-8', 'windows-1252', 'iso-8859-1']  # Lista di codifiche comuni
        text_content = None
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text_content = file.read()
                break  # Se la lettura riesce, interrompe il ciclo
            except UnicodeDecodeError:
                continue  # Prova la prossima codifica
            except IOError as e:
                QMessageBox.critical(self, "Errore di lettura file", f"Impossibile leggere il file: {str(e)}")
                return
            except Exception as e:
                QMessageBox.critical(self, "Errore", f"Errore non previsto: {str(e)}")
                return

        if text_content is None:
            QMessageBox.critical(self, "Errore di lettura",
                                 "Non è stato possibile decodificare il file con le codifiche standard.")
            return

        # Chiedi all'utente dove salvare la presentazione PowerPoint
        output_file, _ = QFileDialog.getSaveFileName(self, "Salva Presentazione", "",
                                                     "PowerPoint Presentation (*.pptx)")
        if not output_file:
            QMessageBox.warning(self, "Attenzione", "Salvataggio annullato. Nessun file selezionato.")
            return

        # Crea la presentazione PowerPoint utilizzando il testo letto dal file
        self.createPresentationFromText(text_content, output_file)

    def generaPresentationConTestoAttuale(self):
        testo_attuale = self.transcriptionTextArea.toPlainText()
        if testo_attuale.strip() == "":
            # Se la QTextEdit è vuota, chiede all'utente di selezionare un file
            file_path, _ = QFileDialog.getOpenFileName(self, "Seleziona File di Testo", "", "Text Files (*.txt)")
            if file_path:
                self.createPresentationFromFile(file_path)
            else:
                # Se l'utente non seleziona un file, mostra un messaggio e non fa nulla
                QMessageBox.warning(self, "Attenzione", "Nessun testo inserito e nessun file selezionato.")
        else:
            # Prompt the user to choose a location to save the presentation
            save_path, _ = QFileDialog.getSaveFileName(self, "Salva Presentazione", "",
                                                       "PowerPoint Presentation (*.pptx)")
            if save_path:
                # Utilizza il testo presente per generare la presentazione e salvare al percorso specificato
                self.createPresentationFromText(testo_attuale, save_path)
            else:
                QMessageBox.warning(self, "Attenzione", "Salvataggio annullato. Nessun file selezionato.")

    def createPresentationFromText(self, testo, output_file):

        # Create a PowerPoint presentation
        prs = Presentation()

        # Select the 'Title and Content' layout which is commonly layout index 1
        title_and_content_layout = prs.slide_layouts[1]

        def imposta_testo_e_font(paragraph, text, size_pt, bold=False):
            """
            Helper function to set text and font properties for a given paragraph.
            """
            # Remove asterisks from the text before setting it
            text = text.replace('*', '')  # Removes all asterisks

            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(size_pt)
            run.font.bold = bold

        # Clean the text: remove format specific asterisks and adjust bullet points
        clean_text = re.sub(r'\*\*(Titolo|Sottotitolo|Contenuto):', r'\1:', testo)  # Remove asterisks around titles
        clean_text = re.sub(r'-\s*', '\u2022 ', clean_text)  # Replace dashes before bullets with bullet points

        # Regex to extract structured information such as title, subtitle, and content
        pattern = r"Titolo:\s*(.*?)\s+Sottotitolo:\s*(.*?)\s+Contenuto:\s*(.*?)\s*(?=Titolo|$)"
        slides_data = re.findall(pattern, clean_text, re.DOTALL)

        for titolo_text, sottotitolo_text, contenuto_text in slides_data:
            # Add a slide with the predefined layout
            slide = prs.slides.add_slide(title_and_content_layout)

            # Set the main title
            titolo = slide.shapes.title
            imposta_testo_e_font(titolo.text_frame.add_paragraph(), titolo_text.strip(), 32, bold=True)

            # Create a textbox for the subtitle directly below the title
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(1)
            sottotitolo_shape = slide.shapes.add_textbox(left, top, width, height)
            sottotitolo_frame = sottotitolo_shape.text_frame
            imposta_testo_e_font(sottotitolo_frame.add_paragraph(), sottotitolo_text.strip(), 24, bold=False)

            # Set the content
            contenuto_box = slide.placeholders[1]
            for line in contenuto_text.strip().split('\n'):
                p = contenuto_box.text_frame.add_paragraph()
                if ':' in line:
                    part1, part2 = line.split(':', 1)
                    imposta_testo_e_font(p, part1.strip() + ':', 20, bold=True)  # Bold the part before the colon
                    imposta_testo_e_font(p, part2.strip(), 20, bold=False)  # Normal text for the part after the colon
                else:
                    imposta_testo_e_font(p, line.strip(), 20, bold=False)  # Normal text if no colon is present

        # Save the presentation if slides have been created
        if prs.slides:
            prs.save(output_file)
            QMessageBox.information(self, "Successo",
                                    "Presentazione PowerPoint generata con successo e salvata in: " + output_file)
        else:
            QMessageBox.warning(self, "Attenzione",
                                "Non sono state generate slides a causa di dati di input non validi o mancanti.")

class CustomComboBox(QComboBox):
    popupOpened = pyqtSignal()

    def showPopup(self):
        self.popupOpened.emit()  # Emetti il segnale quando il popup viene aperto
        super().showPopup()  # Chiamata al metodo originale per assicurarsi che il popup venga mostrato


class CustomTextEdit(QTextEdit):
    def __init__(self, parent=None):
        super().__init__(parent)

    def insertFromMimeData(self, source):
        if source.hasText():
            # Ottieni il testo puro senza formattazione
            plain_text = source.text()
            # Inserisci il testo come testo puro
            self.insertPlainText(plain_text)
        else:
            super().insertFromMimeData(source)
if __name__ == '__main__':
    app = QApplication(sys.argv)
    window = VideoAudioManager()
    window.show()
    sys.exit(app.exec())
